import os
import re
import json
import uuid
import time
import threading
import logging
import requests
import sqlite3
from PyPDF2 import PdfReader
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from flask import Flask, request, jsonify, render_template, g, session
from flask_cors import CORS
from contextlib import closing
from collections import defaultdict

# 全局变量
topology_results = {}
current_upload_filename = None

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("knowledge_graph.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("KnowledgeGraphGenerator")

# 初始化Flask应用
app = Flask(__name__, static_folder='static', template_folder='templates')
CORS(app, resources={r"/*": {"origins": "*"}})  # 增强CORS配置

# 配置会话
app.secret_key = os.environ.get('SECRET_KEY', 'dev_key_for_testing')  # 生产环境应使用环境变量
app.config['SESSION_TYPE'] = 'filesystem'
app.config['SESSION_PERMANENT'] = True
app.config['PERMANENT_SESSION_LIFETIME'] = 86400  # 会话有效期（秒）

# 配置上传文件夹
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# DeepSeek API配置（示例密钥，需替换为实际密钥）
OPENAI_API_KEY = "sk-ba9cc9a26b8c4859ba5c9bad33785093"
MAX_RETRIES = 3
BACKOFF_FACTOR = 2

# 数据库配置
DATABASE = os.path.join(app.root_path, 'knowledge_graph.db')

def get_db():
    """获取数据库连接"""
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
        db.row_factory = sqlite3.Row
    return db

def init_db():
    """初始化数据库（包含问答会话表）"""
    logger.info("开始初始化数据库...")
    with app.app_context():
        db_path = os.path.abspath(DATABASE)
        logger.info(f"数据库文件路径: {db_path}")
        
        # 确保数据库目录存在
        db_dir = os.path.dirname(db_path)
        if not os.path.exists(db_dir):
            os.makedirs(db_dir)
        
        # 如果数据库已存在，跳过初始化
        if os.path.exists(db_path):
            logger.info("数据库已存在，跳过初始化")
            return
        
        logger.info("数据库文件不存在，创建新数据库...")
        with closing(get_db()) as db:
            try:
                schema = """
                CREATE TABLE IF NOT EXISTS topologies (
                    id TEXT PRIMARY KEY,
                    content TEXT NOT NULL,
                    max_nodes INTEGER DEFAULT 0,
                    created_at TEXT
                );

                CREATE TABLE IF NOT EXISTS nodes (
                    id TEXT,
                    topology_id TEXT,
                    label TEXT,
                    level INTEGER,
                    value INTEGER,
                    mastered INTEGER DEFAULT 0,
                    mastery_score REAL DEFAULT 0,
                    consecutive_correct INTEGER DEFAULT 0,
                    content_snippet TEXT,
                    PRIMARY KEY (topology_id, id),
                    FOREIGN KEY (topology_id) REFERENCES topologies (id)
                );

                CREATE TABLE IF NOT EXISTS edges (
                    topology_id TEXT,
                    from_node TEXT,
                    to_node TEXT,
                    label TEXT,
                    PRIMARY KEY (topology_id, from_node, to_node),
                    FOREIGN KEY (topology_id) REFERENCES topologies (id),
                    FOREIGN KEY (from_node) REFERENCES nodes (id),
                    FOREIGN KEY (to_node) REFERENCES nodes (id)
                );

                CREATE TABLE IF NOT EXISTS questions (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    node_id TEXT,
                    question TEXT,
                    answer TEXT,
                    feedback TEXT,
                    correctness INTEGER DEFAULT 0,
                    created_at TEXT,
                    answered_at TEXT,
                    session_id TEXT,
                    FOREIGN KEY (topology_id, node_id) REFERENCES nodes (topology_id, id)
                );

                -- 问答会话表
                CREATE TABLE IF NOT EXISTS quiz_sessions (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    node_id TEXT,
                    created_at TEXT,
                    questions_answered INTEGER DEFAULT 0,
                    consecutive_correct INTEGER DEFAULT 0,
                    mastered INTEGER DEFAULT 0,
                    FOREIGN KEY (topology_id, node_id) REFERENCES nodes (topology_id, id)
                );
                
                -- 历史记录表
                CREATE TABLE IF NOT EXISTS history_records (
                    id TEXT PRIMARY KEY,
                    title TEXT NOT NULL,
                    content TEXT NOT NULL,
                    created_at TEXT NOT NULL,
                    description TEXT,
                    file_path TEXT,
                    user_id TEXT
                );
                """
                db.cursor().executescript(schema)
                db.commit()
                logger.info("数据库表创建成功")
            except Exception as e:
                logger.error(f"数据库初始化失败: {str(e)}", exc_info=True)
                raise

def clean_json_string(s: str) -> str:
    """清洗模型输出，去除Markdown代码块标记"""
    if s is None:
        return ""
    s = re.sub(r"```(?:json)?", "", s)
    return s.strip()

def enhance_json_format(json_str: str) -> str:
    """增强JSON格式，处理各种复杂格式问题"""
    import json
    import re
    
    logger.info(f"原始JSON内容: {json_str[:200]}...")
    
    # 1. 移除可能的前导/尾随非JSON字符
    json_str = re.sub(r'^.*?(\[.*\]).*$', r'\1', json_str, flags=re.DOTALL)
    
    # 2. 处理行首缩进和换行
    json_str = re.sub(r'\n\s*', ' ', json_str)
    
    # 3. 处理可能的单引号问题
    json_str = re.sub(r"'", '"', json_str)
    
    # 4. 处理常见的格式错误（尝试修复）
    try:
        # 尝试直接解析
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        logger.error(f"初始解析失败: {e}")
        
        # 错误定位
        error_pos = e.pos
        logger.error(f"错误位置: {error_pos}, 附近内容: {json_str[error_pos-20:error_pos+20]}")
        
        # 5. 智能修复常见错误
        # 5.1 处理缺少逗号的情况
        if e.msg == "Expecting ',' delimiter":
            logger.info("尝试修复缺少逗号的问题...")
            json_list = list(json_str)
            
            # 在错误位置前查找可能缺少逗号的位置
            search_start = max(0, error_pos - 100)
            bracket_count = 0
            for i in range(error_pos-1, search_start, -1):
                if json_list[i] == ']':
                    bracket_count += 1
                elif json_list[i] == '[':
                    bracket_count -= 1
                
                # 找到合适的位置插入逗号
                if bracket_count == 0 and json_list[i] == ']':
                    json_list.insert(i+1, ',')
                    logger.info(f"在位置 {i+1} 插入逗号")
                    break
            
            # 尝试再次解析
            try:
                return json.loads(''.join(json_list))
            except json.JSONDecodeError as e2:
                logger.error(f"修复后仍失败: {e2}")
        
        # 5.2 处理未闭合的引号
        if e.msg.startswith('Expecting property name enclosed in double quotes'):
            logger.info("尝试修复未闭合的引号...")
            # 查找最近的未闭合引号并添加
            unclosed_quote_pos = json_str.rfind('"', 0, error_pos)
            if unclosed_quote_pos != -1:
                json_list = list(json_str)
                json_list.insert(error_pos, '"')
                logger.info(f"在位置 {error_pos} 添加引号")
                try:
                    return json.loads(''.join(json_list))
                except json.JSONDecodeError as e3:
                    logger.error(f"修复后仍失败: {e3}")
        
        # 6. 使用更宽松的解析器（作为最后的手段）
        try:
            import ast
            logger.info("尝试使用ast.literal_eval进行宽松解析...")
            parsed = ast.literal_eval(json_str)
            # 转换为标准JSON
            return json.loads(json.dumps(parsed))
        except Exception as e4:
            logger.error(f"宽松解析失败: {e4}")
            logger.error(f"无法修复的JSON内容: {json_str}")
            raise RuntimeError("无法解析API返回的JSON格式") from e4

def sanitize_text(text: str) -> str:
    """清理文本中的特殊字符，防止破坏JSON格式"""
    # 移除可能干扰JSON解析的字符
    text = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', text)  # 移除控制字符
    # 转义特殊字符
    text = text.replace('\\', '\\\\')
    text = text.replace('"', '\\"')
    return text

def extract_knowledge_from_text(text: str, max_nodes: int = 0, max_retries: int = MAX_RETRIES) -> list:
    """调用DeepSeek API提取适合树形结构的知识点层级关系"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 清理文本
    sanitized_text = sanitize_text(text)
    
    # 根据节点数量限制调整提示
    node_limit_prompt = ""
    if max_nodes > 0:
        node_limit_prompt = f"请限制生成的知识点数量不超过{max_nodes}个，优先选择最重要的概念。"
    
    # 构建提示
    prompt = f"""
    请分析以下内容，提取其中的知识点及其关系，以便构建一个知识图谱。
    
    {node_limit_prompt}
    
    请按照以下JSON格式返回结果，每个关系表示为一个包含四个元素的数组：
    [
      ["概念A", "关系", "概念B", false],
      ["概念A", "关系", "概念C", false],
      ...
    ]
    
    数组中的第四个元素是布尔值，表示该知识点是否已被掌握（默认为false）。
    关系应该是动词或短语，如"包含"、"属于"、"导致"等。
    请确保概念间的关系是有意义的，并且可以形成一个连贯的知识网络。
    
    以下是需要分析的内容：
    {sanitized_text}
    """
    
    retry_count = 0
    while retry_count < max_retries:
        try:
            logger.info(f"尝试调用API提取知识点关系，第{retry_count+1}次尝试")
            
            # 创建符合OpenAI API要求的消息格式
            messages = [
                {"role": "system", "content": "你是一个专业的知识图谱构建助手，擅长从文本中提取概念及其关系。"},
                {"role": "user", "content": prompt}
            ]
            
            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=messages,
                temperature=0.2,
                max_tokens=4000
            )
            
            result = response.choices[0].message.content if response.choices and len(response.choices) > 0 else ""
            if result:
                logger.info(f"API返回原始结果: {result[:200]}...")
                
                # 清理结果
                cleaned_json = clean_json_string(result)
                
                # 增强JSON格式处理
                try:
                    # 解析JSON
                    knowledge_edges = enhance_json_format(cleaned_json)
                    
                    # 检查是否已经是新格式（包含四个元素的数组或对象格式）
                    enhanced_knowledge_edges = []
                    for edge in knowledge_edges:
                        if isinstance(edge, list):
                            if len(edge) >= 4:
                                # 已经是四元组格式，直接转换为对象格式
                                enhanced_edge = {
                                    "source": edge[0],
                                    "relation": edge[1],
                                    "target": edge[2],
                                    "highlighted": edge[3] if isinstance(edge[3], bool) else False
                                }
                            elif len(edge) >= 3:
                                # 三元组格式，添加highlighted字段
                                enhanced_edge = {
                                    "source": edge[0],
                                    "relation": edge[1],
                                    "target": edge[2],
                                    "highlighted": False  # 默认未点亮
                                }
                            else:
                                logger.warning(f"跳过格式不正确的边: {edge}")
                                continue
                        elif isinstance(edge, dict) and "source" in edge and "relation" in edge and "target" in edge:
                            # 已经是对象格式，确保有highlighted字段
                            if "highlighted" not in edge:
                                edge["highlighted"] = False
                            enhanced_edge = edge
                        else:
                            logger.warning(f"跳过格式不正确的边: {edge}")
                            continue
                        
                        enhanced_knowledge_edges.append(enhanced_edge)
                    
                    logger.info(f"成功提取知识点关系，共{len(enhanced_knowledge_edges)}条")
                    
                    # 保存API响应到文件（调试用）
                    timestamp = time.strftime('%Y%m%d_%H%M%S')
                    response_file = f"api_response_{timestamp}.txt"
                    # 确保使用绝对路径保存文件
                    response_file_path = os.path.join(os.getcwd(), response_file)
                    
                    with open(response_file_path, "w", encoding="utf-8") as f:
                        json.dump(enhanced_knowledge_edges, f, ensure_ascii=False, indent=2)
                    logger.info(f"API响应已保存到文件: {response_file_path}")
                    
                    # 自动保存到历史记录
                    try:
                        with app.app_context():
                            # 使用上传的文件名作为标题
                            global current_upload_filename
                            title = f"{current_upload_filename} - 知识图谱" if current_upload_filename else f"知识图谱 {timestamp}"
                            history_id = str(uuid.uuid4())
                            created_at = time.strftime('%Y-%m-%d %H:%M:%S')
                            
                            # 保存到数据库
                            with get_db() as db:
                                db.execute(
                                    'INSERT INTO history_records (id, title, content, created_at, description, file_path, user_id) VALUES (?, ?, ?, ?, ?, ?, ?)',
                                    (history_id, title, json.dumps(enhanced_knowledge_edges), created_at, "自动保存的知识图谱", response_file_path, str(uuid.uuid4()))
                                )
                                db.commit()
                            logger.info(f"知识图谱已自动保存到历史记录，ID: {history_id}, 标题: {title}, 文件路径: {response_file_path}")
                    except Exception as e:
                        logger.error(f"保存历史记录失败: {str(e)}", exc_info=True)
                    
                    return enhanced_knowledge_edges
                    
                except Exception as json_err:
                    logger.error(f"JSON解析失败: {str(json_err)}", exc_info=True)
                    retry_count += 1
                    time.sleep(BACKOFF_FACTOR * retry_count)  # 指数退避
            else:
                logger.error("API返回空结果")
                retry_count += 1
                time.sleep(BACKOFF_FACTOR * retry_count)
                
        except Exception as e:
            logger.error(f"API调用失败: {str(e)}", exc_info=True)
            retry_count += 1
            time.sleep(BACKOFF_FACTOR * retry_count)  # 指数退避
    
    logger.error(f"达到最大重试次数({max_retries})，无法提取知识点关系")
    raise RuntimeError(f"无法从文本中提取知识点关系，请检查API配置或重试")

def parse_document(file_path):
    """解析不同类型的文档，提取文本内容"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_ext == '.pdf':
            return parse_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            return parse_docx(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return parse_pptx(file_path)
        elif file_ext == '.txt':
            return parse_txt(file_path)
        elif file_ext in ['.html', '.htm']:
            return parse_html(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {file_ext}")
    except Exception as e:
        logger.error(f"解析文档失败: {str(e)}", exc_info=True)
        raise

def parse_pdf(file_path):
    """解析PDF文件"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                text += page_text + "\n\n"
                if i % 10 == 0:
                    logger.info(f"已解析PDF第 {i} 页")
    except Exception as e:
        logger.error(f"PDF解析错误: {str(e)}", exc_info=True)
        raise
    return text

def parse_docx(file_path):
    """解析Word文档"""
    text = ""
    try:
        doc = Document(file_path)
        for i, para in enumerate(doc.paragraphs):
            if para.text:
                text += para.text + "\n"
            if i % 50 == 0:
                logger.info(f"已解析Word第 {i} 段落")
    except Exception as e:
        logger.error(f"Word解析错误: {str(e)}", exc_info=True)
        raise
    return text

def parse_pptx(file_path):
    """解析PowerPoint文档"""
    text = ""
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # 检查是否有text属性并且可以安全访问
                if hasattr(shape, "text_frame") and shape.text_frame:
                    if hasattr(shape.text_frame, "text") and shape.text_frame.text:
                        text += shape.text_frame.text + "\n"
            text += "\n"
            if i % 10 == 0:
                logger.info(f"已解析PPT第 {i} 张幻灯片")
    except Exception as e:
        logger.error(f"PowerPoint解析错误: {str(e)}", exc_info=True)
        raise
    return text

def parse_txt(file_path):
    """解析纯文本文件"""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def parse_html(file_path):
    """解析HTML文件"""
    with open(file_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    soup = BeautifulSoup(html_content, 'lxml')
    text = soup.get_text()
    return ' '.join(text.split())

def extract_content_snippet(content: str, topic: str) -> str:
    """从原文中提取与主题相关的片段"""
    # 查找主题在内容中的位置
    index = content.lower().find(topic.lower())
    if index == -1:
        return ""
    
    # 提取上下文片段（主题前后各200个字符）
    start = max(0, index - 200)
    end = min(len(content), index + len(topic) + 200)
    snippet = content[start:end]
    
    # 确保片段包含主题
    if topic.lower() not in snippet.lower():
        return ""
    
    # 添加省略号表示截断
    if start > 0:
        snippet = "..." + snippet
    if end < len(content):
        snippet = snippet + "..."
    
    return snippet

def build_tree_structure(knowledge_edges, topology_id, content: str, max_nodes: int = 0):
    """构建树形结构，计算节点层级和重要性"""
    logger.info(f"开始构建树形结构，共{len(knowledge_edges)}条关系")
    
    nodes = {}
    edges = []
    all_node_ids = set()
    
    # 处理每条边，创建节点和边的数据结构
    for edge in knowledge_edges:
        # 检查是否为新格式（对象格式）
        if isinstance(edge, dict) and 'source' in edge and 'target' in edge and 'relation' in edge:
            src = edge['source']
            rel = edge['relation']
            tgt = edge['target']
            highlighted = False  # 强制设置为False，新生成的图谱节点都不应该被点亮
        elif isinstance(edge, list) and len(edge) >= 3:
            # 旧格式（三元组数组）
            src, rel, tgt = edge
            highlighted = False
        else:
            logger.warning(f"跳过格式不正确的边: {edge}")
            continue
        
        all_node_ids.add(src)
        all_node_ids.add(tgt)
            
        # 创建或更新源节点
        if src not in nodes:
            nodes[src] = {
                "id": src,
                "label": src,
                "level": 0,  # 初始层级，后续会更新
                "value": 1,  # 初始重要性，后续会更新
                "mastered": 0,
                "mastery_score": 0,
                "consecutive_correct": 0,
                "content_snippet": extract_content_snippet(content, src),
                "highlighted": highlighted
            }
        
        # 创建或更新目标节点
        if tgt not in nodes:
            nodes[tgt] = {
                "id": tgt,
                "label": tgt,
                "level": 0,  # 初始层级，后续会更新
                "value": 1,  # 初始重要性，后续会更新
                "mastered": 0,
                "mastery_score": 0,
                "consecutive_correct": 0,
                "content_snippet": extract_content_snippet(content, tgt),
                "highlighted": highlighted
            }
        
        # 创建边
        edges.append({
            "from_node": src,
            "to_node": tgt,
            "label": rel
        })
    
    # 计算节点层级
    def calculate_level(node_id, current_level=0):
        if node_id in nodes:
            nodes[node_id]["level"] = max(nodes[node_id]["level"], current_level)
            # 递归设置子节点层级
            for edge in edges:
                if edge["from_node"] == node_id:
                    calculate_level(edge["to_node"], current_level + 1)
    
    # 找到根节点（没有父节点的节点）- 使用原版逻辑
    root_candidates = all_node_ids.copy()
    for edge in knowledge_edges:
        # 兼容新旧格式
        if isinstance(edge, dict) and 'target' in edge:
            tgt = edge['target']
        elif isinstance(edge, list) and len(edge) >= 3:
            tgt = edge[2]
        else:
            continue
        root_candidates.discard(tgt)
    
    root = next(iter(root_candidates)) if root_candidates else (list(nodes.keys())[0] if nodes else None)
    
    # 从根节点开始计算层级
    if root:
        calculate_level(root, 0)
    
    # 计算节点重要性（基于连接数和层级）
    for node_id in nodes:
        # 连接数（入度+出度）
        in_connections = sum(1 for edge in edges if edge["to_node"] == node_id)
        out_connections = sum(1 for edge in edges if edge["from_node"] == node_id)
        connection_count = in_connections + out_connections
        
        # 层级因子（层级越低越重要）
        level_factor = 1.0 / (nodes[node_id]["level"] + 1)
        # 重要性得分
        nodes[node_id]["value"] = max(1, int(connection_count * level_factor * 5))
    
    # 如果设置了最大节点数，进行过滤
    if max_nodes > 0 and len(nodes) > max_nodes:
        # 按重要性排序
        sorted_nodes = sorted(nodes.items(), key=lambda x: x[1]["value"], reverse=True)
        # 保留最重要的节点
        kept_nodes = dict(sorted_nodes[:max_nodes])
        # 过滤边
        kept_edges = [edge for edge in edges 
                     if edge["from_node"] in kept_nodes and edge["to_node"] in kept_nodes]
        
        nodes = kept_nodes
        edges = kept_edges
    
    logger.info(f"树形结构构建完成，共{len(nodes)}个节点，{len(edges)}条边")
    return {"nodes": list(nodes.values()), "edges": edges}

def save_to_database(topology_id, nodes, edges, content: str, max_nodes=0):
    """将知识图谱数据保存到数据库（保存原文内容和节点数量限制）"""
    with app.app_context():
        db = get_db()
        cursor = db.cursor()
        
        # 保存拓扑图信息（包含原文内容和节点数量限制）
        cursor.execute(
            "INSERT OR REPLACE INTO topologies (id, content, max_nodes, created_at) VALUES (?, ?, ?, ?)",
            (topology_id, content, max_nodes, time.strftime('%Y-%m-%d %H:%M:%S'))
        )
        
        # 保存节点
        for node in nodes:
            cursor.execute(
                """INSERT OR REPLACE INTO nodes 
                (topology_id, id, label, level, value, mastered, mastery_score, consecutive_correct, content_snippet) 
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (topology_id, node["id"], node["label"], node["level"], node["value"], 
                 1 if node.get("highlighted", False) else 0,  # 使用highlighted字段作为mastered的值
                 node.get("mastery_score", 0.0), 
                 node.get("consecutive_correct", 0), 
                 node.get("content_snippet", ""))
            )
        
        # 保存边
        for edge in edges:
            cursor.execute(
                "INSERT OR REPLACE INTO edges (topology_id, from_node, to_node, label) VALUES (?, ?, ?, ?)",
                (topology_id, edge["from_node"], edge["to_node"], edge["label"])
            )
        
        db.commit()

def process_document(file_path, topology_id, max_nodes=0):
    """处理文档，提取知识图谱"""
    try:
        # 更新进度
        update_progress(topology_id, 10, "解析文档内容...")
        
        # 解析文档
        text = parse_document(file_path)
        if not text:
            raise ValueError("无法从文档中提取文本内容")
        
        # 如果文本过长，截取前8000字符
        if len(text) > 8000:
            logger.info("文档内容过长，截取前8000字符")
            text = text[:8000]
        
        # 更新进度
        update_progress(topology_id, 20, "准备提取知识层级...")
        
        # 更新进度
        update_progress(topology_id, 60, "调用DeepSeek API提取知识层级...")
        
        # 提取知识点关系
        knowledge_edges = extract_knowledge_from_text(text, max_nodes)
        if not knowledge_edges:
            raise ValueError("未能从文本中提取知识层级关系")
        
        logger.info(f"成功提取{len(knowledge_edges)}条知识层级关系")
        
        # 更新进度
        update_progress(topology_id, 80, "构建树形知识图并提取原文片段...")
        
        # 构建树形结构
        knowledge_graph = build_tree_structure(knowledge_edges, topology_id, text, max_nodes)
        
        # 保存到数据库
        save_to_database(topology_id, knowledge_graph["nodes"], knowledge_graph["edges"], text, max_nodes)
        
        # 格式化为前端需要的格式
        formatted_nodes = []
        for node in knowledge_graph["nodes"]:
            formatted_node = {
                "id": node["id"],
                "label": node["label"],
                "level": node["level"],
                "value": node["value"],
                "mastered": bool(node.get("highlighted", False)),
                "mastery_score": node.get("mastery_score", 0.0),
                "consecutive_correct": node.get("consecutive_correct", 0),
                "content_snippet": node.get("content_snippet", "")
            }
            formatted_nodes.append(formatted_node)
        
        formatted_edges = []
        for edge in knowledge_graph["edges"]:
            formatted_edge = {
                "from": edge["from_node"],
                "to": edge["to_node"],
                "label": edge["label"],
                "arrows": "to"
            }
            formatted_edges.append(formatted_edge)
        
        # 构建完整数据
        frontend_data = {
            "nodes": formatted_nodes,
            "edges": formatted_edges,
            "root": knowledge_graph["nodes"][0]["id"] if knowledge_graph["nodes"] else None
        }
        
        # 更新全局变量
        global topology_results
        topology_results[topology_id] = {
            "status": "completed",
            "data": frontend_data,
            "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
            "node_count": len(formatted_nodes),
            "edge_count": len(formatted_edges),
            "max_nodes": max_nodes
        }
        
        # 更新进度
        update_progress(topology_id, 100, "知识图谱生成完成")
        
        return frontend_data
    except Exception as e:
        logger.error(f"处理文档出错: {str(e)}", exc_info=True)
        # 更新进度为错误状态
        update_progress(topology_id, -1, f"处理过程中出错: {str(e)}")
        raise

def update_progress(topology_id, progress, message):
    """更新处理进度"""
    global topology_results
    
    if 'topology_results' not in globals():
        topology_results = {}
    
    if topology_id not in topology_results:
        topology_results[topology_id] = {
            "status": "processing",
            "progress": 0,
            "message": "初始化...",
            "created_at": time.strftime('%Y-%m-%d %H:%M:%S')
        }
    
    # 更新进度信息
    if progress < 0:
        # 错误状态
        topology_results[topology_id]["status"] = "error"
        topology_results[topology_id]["message"] = message
    elif progress >= 100:
        # 完成状态
        topology_results[topology_id]["status"] = "completed"
        topology_results[topology_id]["progress"] = 100
        topology_results[topology_id]["message"] = message
    else:
        # 处理中状态
        topology_results[topology_id]["status"] = "processing"
        topology_results[topology_id]["progress"] = progress
        topology_results[topology_id]["message"] = message
    
    logger.info(f"拓扑ID: {topology_id}, 进度: {progress}%, 消息: {message}")

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/generate', methods=['POST'])
def generate_knowledge_graph():
    """处理用户点击生成按钮后的请求，支持节点数量控制"""
    if 'file' not in request.files:
        logger.error("文件上传错误: 没有文件")
        return jsonify({'status': 'error', 'message': '没有文件上传'}), 400
    
    file = request.files['file']
    if file.filename == '':
        logger.error("文件上传错误: 未选择文件")
        return jsonify({'status': 'error', 'message': '未选择文件'}), 400
    
    # 获取节点数量 - 确保从表单获取
    max_nodes = request.form.get('max_nodes', 0, type=int)
    
    # 检查文件大小
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    
    if file_size > 50 * 1024 * 1024:  # 50MB限制
        logger.error(f"文件上传错误: 文件过大 ({file_size/1024/1024:.2f} MB)")
        return jsonify({'status': 'error', 'message': '文件大小超过50MB限制'}), 400
    
    topology_id = str(uuid.uuid4())
    original_filename = file.filename if file.filename else "未命名文件"
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{topology_id}_{original_filename}")
    file.save(file_path)
    
    logger.info(f"文件上传成功: {file_path}, 大小: {file_size/1024/1024:.2f} MB, 最大节点数: {max_nodes}")
    
    # 保存文件名到全局变量，用于后续生成历史记录标题
    global current_upload_filename
    current_upload_filename = os.path.splitext(original_filename)[0] if original_filename else "未命名文件"
    
    # 启动处理线程，并在应用上下文中执行
    threading.Thread(
        target=lambda: with_app_context(process_document, file_path, topology_id, max_nodes)
    ).start()
    
    return jsonify({
        'status': 'success',
        'topology_id': topology_id,
        'message': '文档上传成功，正在生成知识图谱',
        'max_nodes': max_nodes,  # 返回节点数量限制
        'filename': original_filename  # 返回原始文件名
    })

def with_app_context(func, *args, **kwargs):
    """在应用上下文中执行函数"""
    with app.app_context():
        func(*args, **kwargs)

@app.route('/api/topology/<topology_id>', methods=['GET'])
def get_topology(topology_id):
    """获取知识图谱数据"""
    global topology_results
    
    # 初始化全局变量
    if 'topology_results' not in globals():
        topology_results = {}
    
    try:
        # 检查是否有缓存的处理结果
        if topology_id in topology_results:
            result = topology_results[topology_id]
            
            # 如果处理完成，返回完整数据
            if result.get("status") == "completed":
                return jsonify({
                    "status": "success",
                    "data": result.get("data", {}),
                    "node_count": result.get("node_count", 0),
                    "edge_count": result.get("edge_count", 0)
                })
            
            # 如果处理出错，返回错误信息
            elif result.get("status") == "error":
                return jsonify({
                    "status": "error",
                    "message": result.get("message", "未知错误")
                }), 500
            
            # 如果正在处理中，返回进度信息
            else:
                return jsonify({
                    "status": "processing",
                    "progress": result.get("progress", 0),
                    "message": result.get("message", "正在处理...")
                })
        
        # 如果没有缓存，尝试从数据库获取
        with get_db() as db:
            # 获取拓扑图基本信息
            cursor = db.execute(
                "SELECT id, content, max_nodes, created_at FROM topologies WHERE id = ?",
                (topology_id,)
            )
            topology = cursor.fetchone()
            
            if not topology:
                return jsonify({"status": "error", "message": "找不到指定的知识图谱"}), 404
            
            # 获取节点数据
            cursor = db.execute(
                "SELECT id, label, level, value, mastered, mastery_score, consecutive_correct, content_snippet "
                "FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            nodes = [dict(row) for row in cursor.fetchall()]
            
            # 获取边数据
            cursor = db.execute(
                "SELECT from_node, to_node, label FROM edges WHERE topology_id = ?",
                (topology_id,)
            )
            edges = [dict(row) for row in cursor.fetchall()]
            
            # 格式化为前端需要的格式
            formatted_nodes = []
            for node in nodes:
                formatted_node = {
                    "id": node["id"],
                    "label": node["label"],
                    "level": node["level"],
                    "value": node["value"],
                    "mastered": bool(node["mastered"]),
                    "mastery_score": node["mastery_score"],
                    "consecutive_correct": node["consecutive_correct"],
                    "content_snippet": node["content_snippet"],
                    "highlighted": bool(node["mastered"])  # 使用mastered字段作为highlighted的值
                }
                formatted_nodes.append(formatted_node)
            
            formatted_edges = []
            for edge in edges:
                formatted_edge = {
                    "from": edge["from_node"],
                    "to": edge["to_node"],
                    "label": edge["label"],
                    "arrows": "to"
                }
                formatted_edges.append(formatted_edge)
            
            # 找出根节点
            all_targets = {edge["to"] for edge in formatted_edges}
            root_candidates = {node["id"] for node in formatted_nodes} - all_targets
            root = next(iter(root_candidates)) if root_candidates else (formatted_nodes[0]["id"] if formatted_nodes else None)
            
            # 构建完整数据
            graph_data = {
                "nodes": formatted_nodes,
                "edges": formatted_edges,
                "root": root
            }
            
            # 缓存结果
            topology_results[topology_id] = {
                "status": "completed",
                "data": graph_data,
                "created_at": topology["created_at"],
                "node_count": len(formatted_nodes),
                "edge_count": len(formatted_edges),
                "max_nodes": topology["max_nodes"]
            }
            
            return jsonify({
                "status": "success",
                "data": graph_data,
                "node_count": len(formatted_nodes),
                "edge_count": len(formatted_edges)
            })
            
    except Exception as e:
        logger.error(f"获取拓扑图错误: {str(e)}", exc_info=True)
        return jsonify({"status": "error", "message": str(e)}), 500

@app.route('/api/topology/<topology_id>/set_max_nodes', methods=['POST'])
def set_topology_max_nodes(topology_id):
    """更新拓扑图的节点数量设置"""
    try:
        data = request.json
        max_nodes = data.get('max_nodes', 0)
        
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 更新拓扑图的节点数量设置
            cursor.execute(
                "UPDATE topologies SET max_nodes = ? WHERE id = ?",
                (max_nodes, topology_id)
            )
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'message': '节点数量设置已更新',
                'max_nodes': max_nodes
            })
            
    except Exception as e:
        logger.error(f"设置节点数量错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"设置节点数量时出错: {str(e)}"
        }), 500

@app.route('/api/topology/<topology_id>/regenerate', methods=['POST'])
def regenerate_topology(topology_id):
    """重新生成知识图谱，使用用户输入的新节点数量"""
    try:
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 从请求中获取新的节点数量
            data = request.json
            max_nodes = data.get('max_nodes', 0)  # 从请求中获取新的节点数量
            
            # 从数据库获取原文内容
            cursor.execute(
                "SELECT content FROM topologies WHERE id = ?",
                (topology_id,)
            )
            topology = cursor.fetchone()
            
            if not topology:
                return jsonify({
                    'status': 'error',
                    'message': '拓扑图不存在'
                }), 404
            
            content = topology["content"]
            
            update_progress(topology_id, 30, "重新提取知识层级...")
            knowledge_edges = extract_knowledge_from_text(content, max_nodes)  # 使用新的节点数量
            logger.info(f"重新生成成功提取{len(knowledge_edges)}条知识层级关系")
            
            # 保存当前节点的掌握状态
            cursor.execute(
                "SELECT id, mastered, mastery_score, consecutive_correct FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            mastery_states = {row["id"]: dict(row) for row in cursor.fetchall()}
            
            update_progress(topology_id, 70, "重新构建树形知识图...")
            knowledge_graph = build_tree_structure(knowledge_edges, topology_id, content, max_nodes)  # 使用新的节点数量
            
            # 恢复节点的掌握状态
            for node in knowledge_graph["nodes"]:
                node_id = node["id"]
                if node_id in mastery_states:
                    state = mastery_states[node_id]
                    node["mastered"] = bool(state["mastered"])
                    node["mastery_score"] = state["mastery_score"]
                    node["consecutive_correct"] = state["consecutive_correct"]
                    
                    # 更新数据库中的节点状态
                    cursor.execute(
                        """UPDATE nodes SET 
                        mastered = ?, mastery_score = ?, consecutive_correct = ?
                        WHERE topology_id = ? AND id = ?""",
                        (int(state["mastered"]), state["mastery_score"], 
                         state["consecutive_correct"], topology_id, node_id)
                    )
            
            # 更新拓扑图的节点数量设置到数据库
            cursor.execute(
                "UPDATE topologies SET max_nodes = ? WHERE id = ?",
                (max_nodes, topology_id)
            )
            
            # 更新处理结果
            topology_results[topology_id] = {
                "status": "completed",
                "data": knowledge_graph,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph["nodes"]),
                "edge_count": len(knowledge_graph["edges"]),
                "processing_time": 0,
                "text_length": len(content),
                "max_nodes": max_nodes  # 保存新的节点数量限制
            }
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'message': '知识图谱重新生成成功',
                'node_count': len(knowledge_graph["nodes"]),
                'edge_count': len(knowledge_graph["edges"]),
                'max_nodes': max_nodes  # 返回新的节点数量限制
            })
            
    except Exception as e:
        logger.error(f"重新生成知识图谱错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"重新生成知识图谱时出错: {str(e)}"
        }), 500


@app.route('/api/topology/<topology_id>/node/<node_id>/question', methods=['GET'])
def get_question(topology_id, node_id):
    """获取关于指定节点的问题（基于原文内容，支持会话管理）"""
    try:
        with app.app_context():
            # 从数据库获取节点信息
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            
            if not node:
                return jsonify({
                    'status': 'error',
                    'message': '节点不存在'
                }), 404
            
            node_label = node["label"]
            content_snippet = node["content_snippet"]
            
            # 检查是否已有活跃会话
            session_id = request.args.get('session_id')
            if session_id:
                cursor.execute(
                    "SELECT mastered, consecutive_correct FROM quiz_sessions WHERE id = ?",
                    (session_id,)
                )
                session = cursor.fetchone()
                if session and session["mastered"]:
                    return jsonify({
                        'status': 'success',
                        'mastered': True,
                        'message': '该知识点已掌握'
                    })
            
            # 创建或获取会话
            if not session_id:
                session_id = str(uuid.uuid4())
                cursor.execute(
                    "INSERT INTO quiz_sessions (id, topology_id, node_id, created_at) VALUES (?, ?, ?, ?)",
                    (session_id, topology_id, node_id, time.strftime('%Y-%m-%d %H:%M:%S'))
                )
                db.commit()
            
            # 获取会话状态
            cursor.execute(
                "SELECT consecutive_correct FROM quiz_sessions WHERE id = ?",
                (session_id,)
            )
            session = cursor.fetchone()
            consecutive_correct = session["consecutive_correct"] if session else 0
            
            # 生成问题（基于会话状态）
            question = generate_question(node_label, content_snippet, consecutive_correct)
            
            # 保存问题到数据库
            question_id = str(uuid.uuid4())
            cursor.execute(
                "INSERT INTO questions (id, topology_id, node_id, question, session_id, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                (question_id, topology_id, node_id, question, session_id, time.strftime('%Y-%m-%d %H:%M:%S'))
            )
            db.commit()
            
            return jsonify({
                'status': 'success',
                'data': {
                    'question_id': question_id,
                    'question': question,
                    'node_id': node_id,
                    'session_id': session_id
                }
            })
    except Exception as e:
        logger.error(f"获取问题错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"生成问题时出错: {str(e)}"
        }), 500

def generate_question(topic, context, consecutive_correct=0):
    """根据连续正确次数生成不同难度的问题"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 根据掌握程度生成不同难度的问题
    difficulty_map = {
        0: "基础概念题，用简洁的语言解释",
        1: "理解应用题，结合实例说明",
        2: "综合分析题，比较相关概念"
    }
    difficulty = difficulty_map.get(consecutive_correct, "进阶思考题，拓展应用场景")
    
    # 构建提示词
    messages = [
        {"role": "system", "content": f"""你是一个教育专家，能够基于原文内容生成有针对性的问题。
请生成一个{difficulty}，测试用户对"{topic}"的理解。
问题应该清晰明确，基于提供的原文内容，难度适合当前掌握程度。
只返回问题文本，不需要其他内容。"""},
        {"role": "user", "content": f"原文片段: {context}\n问题:"}
    ]
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=150
        )
        
        question = response.choices[0].message.content.strip()
        return question
    except Exception as e:
        logger.error(f"生成问题出错: {str(e)}", exc_info=True)
        return f"关于{topic}的问题（基于原文）"

@app.route('/api/topology/<topology_id>/question/<question_id>/answer', methods=['POST'])
def answer_question(topology_id, question_id):
    """处理用户对问题的回答（支持会话状态管理）并更新节点状态"""
    try:
        with app.app_context():
            data = request.json
            answer = data.get('answer', '')
            node_id = data.get('node_id', '')
            session_id = data.get('session_id', '')
            
            if not answer or not node_id or not session_id:
                return jsonify({
                    'status': 'error',
                    'message': '缺少必要的参数'
                }), 400
            
            # 从数据库获取问题
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT question, node_id, session_id FROM questions WHERE id = ? AND topology_id = ?",
                (question_id, topology_id)
            )
            question_data = cursor.fetchone()
            
            if not question_data:
                return jsonify({
                    'status': 'error',
                    'message': '问题不存在'
                }), 404
            
            question = question_data["question"]
            stored_node_id = question_data["node_id"]
            stored_session_id = question_data["session_id"]
            
            if stored_node_id != node_id or stored_session_id != session_id:
                return jsonify({
                    'status': 'error',
                    'message': '问题与会话不匹配'
                }), 400
            
            # 从数据库获取节点信息
            cursor.execute(
                "SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            if not node:
                return jsonify({
                    'status': 'error',
                    'message': '节点不存在'
                }), 400
            
            node_label = node["label"]
            content_snippet = node["content_snippet"]
            
            # 调用DeepSeek评估回答
            evaluation = evaluate_answer(question, answer, node_label, content_snippet)
            
            # 确定回答是否正确
            is_correct = evaluation["correct"] if "correct" in evaluation else False
            feedback_text = evaluation["feedback"] if "feedback" in evaluation else "无法评估回答"
            
            # 更新问题状态
            cursor.execute(
                "UPDATE questions SET answered_at = ?, answer = ?, feedback = ?, correctness = ? WHERE id = ?",
                (time.strftime('%Y-%m-%d %H:%M:%S'), answer, feedback_text, 1 if is_correct else 0, question_id)
            )
            
            # 更新会话状态
            cursor.execute(
                "SELECT consecutive_correct, mastered FROM quiz_sessions WHERE id = ?",
                (session_id,)
            )
            session = cursor.fetchone()
            if not session:
                return jsonify({
                    'status': 'error',
                    'message': '问答会话不存在'
                }), 404
            
            current_consecutive = session["consecutive_correct"]
            current_mastered = session["mastered"]
            
            # 更新连续正确计数
            new_consecutive = current_consecutive + 1 if is_correct else 0
            new_mastered = 1 if new_consecutive >= 3 else 0
            
            cursor.execute(
                """UPDATE quiz_sessions SET 
                questions_answered = questions_answered + 1,
                consecutive_correct = ?,
                mastered = ?
                WHERE id = ?""",
                (new_consecutive, new_mastered, session_id)
            )
            
            # 更新节点的掌握状态
            cursor.execute(
                "SELECT mastery_score, consecutive_correct FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node_status = cursor.fetchone()
            current_node_score = node_status["mastery_score"] if node_status else 0
            current_node_consecutive = node_status["consecutive_correct"] if node_status else 0
            
            node_new_score = min(10, current_node_score + (1 if is_correct else -0.5))
            node_new_consecutive = new_consecutive
            node_new_mastered = new_mastered
            
            cursor.execute(
                "UPDATE nodes SET mastery_score = ?, consecutive_correct = ?, mastered = ? WHERE topology_id = ? AND id = ?",
                (node_new_score, node_new_consecutive, int(node_new_mastered), topology_id, node_id)
            )
            
            # 确保掌握状态正确更新（显式处理）
            if node_new_mastered:
                cursor.execute(
                    "UPDATE nodes SET mastered = 1 WHERE topology_id = ? AND id = ?",
                    (topology_id, node_id)
                )
            
            # 如果未掌握，生成下一个问题
            next_question = None
            next_question_id = None
            if not new_mastered:
                # 生成下一个问题（基于更新后的状态）
                next_question = generate_question(node_label, content_snippet, new_consecutive)
                if next_question:
                    next_question_id = str(uuid.uuid4())
                    cursor.execute(
                        "INSERT INTO questions (id, topology_id, node_id, question, session_id, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                        (next_question_id, topology_id, node_id, next_question, session_id, time.strftime('%Y-%m-%d %H:%M:%S'))
                    )
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'data': {
                    'correct': is_correct,
                    'feedback': feedback_text,
                    'mastered': new_mastered,
                    'consecutive_correct': new_consecutive,
                    'session_id': session_id,
                    'next_question': {
                        'id': next_question_id,
                        'question': next_question
                    } if next_question else None
                }
            })
    except Exception as e:
        logger.error(f"处理回答错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"处理回答时出错: {str(e)}"
        }), 500

def evaluate_answer(question, answer, topic, context):
    """评估用户回答，使用DeepSeek API进行评估"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 构建评估提示
    prompt = f"""
    你是一个教育评估专家，需要评估学生对知识点的理解程度。
    
    知识点: {topic}
    相关上下文: {context}
    
    问题: {question}
    学生回答: {answer}
    
    请评估学生的回答是否正确，并给出评分（0-10分）和具体反馈。
    
    请按照以下JSON格式返回结果:
    {{
        "score": 评分（0-10的整数）,
        "correct": true/false（是否基本正确）,
        "feedback": "详细的评价反馈"
    }}
    
    注意:
    1. 评分标准: 0-3分表示完全不理解，4-6分表示部分理解，7-10分表示良好理解
    2. 即使答案不完全正确，如果展示了对核心概念的理解，也应给予一定分数
    3. 反馈应该具体指出答案的优点和不足，并给出改进建议
    """
    
    try:
        # 创建符合OpenAI API要求的消息格式
        messages = [
            {"role": "system", "content": "你是一个教育评估专家，负责评估学生对知识点的理解程度。"},
            {"role": "user", "content": prompt}
        ]
        
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            temperature=0.2,
            max_tokens=1000
        )
        
        result = response.choices[0].message.content if response.choices and len(response.choices) > 0 else ""
        if not result:
            logger.error("评估API返回空结果")
            return {"score": 5, "correct": False, "feedback": "系统无法评估您的回答，请重试。"}
        
        # 清理结果
        cleaned_json = clean_json_string(result)
        
        try:
            # 解析JSON
            evaluation = json.loads(cleaned_json)
            
            # 验证必要字段
            if "score" not in evaluation or "correct" not in evaluation or "feedback" not in evaluation:
                logger.error(f"评估结果缺少必要字段: {evaluation}")
                return {"score": 5, "correct": False, "feedback": "系统评估结果格式错误，请重试。"}
            
            # 确保分数在有效范围内
            evaluation["score"] = max(0, min(10, int(evaluation.get("score", 5))))
            
            return evaluation
            
        except json.JSONDecodeError as e:
            logger.error(f"评估结果JSON解析失败: {str(e)}")
            return {"score": 5, "correct": False, "feedback": "系统无法解析评估结果，请重试。"}
            
    except Exception as e:
        logger.error(f"评估API调用失败: {str(e)}", exc_info=True)
        return {"score": 5, "correct": False, "feedback": "评估过程中出错，请重试。"}

@app.route('/api/topology/<topology_id>/ignore', methods=['POST'])
def ignore_nodes(topology_id):
    """忽略用户选择的节点"""
    try:
        with app.app_context():
            data = request.json
            ignored_nodes = data.get('ignored_nodes', [])
            
            db = get_db()
            cursor = db.cursor()
            
            # 获取所有节点
            cursor.execute(
                "SELECT id, label, level, value, mastered, mastery_score, content_snippet FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            all_nodes = [dict(row) for row in cursor.fetchall()]
            
            # 获取所有边
            cursor.execute(
                "SELECT from_node, to_node, label FROM edges WHERE topology_id = ?",
                (topology_id,)
            )
            all_edges = [dict(row) for row in cursor.fetchall()]
            
            # 筛选节点（排除被忽略的）
            filtered_nodes = [node for node in all_nodes if node["id"] not in ignored_nodes]
            
            # 筛选边（只保留未被忽略节点之间的边）
            filtered_edges = []
            node_ids = set([node["id"] for node in filtered_nodes])
            for edge in all_edges:
                if edge["from_node"] in node_ids and edge["to_node"] in node_ids:
                    filtered_edges.append(edge)
            
            return jsonify({
                'status': 'success',
                'data': {
                    'nodes': filtered_nodes,
                    'edges': filtered_edges,
                    'root': next((node["id"] for node in filtered_nodes if node["level"] == 0), filtered_nodes[0]["id"] if filtered_nodes else None)
                }
            })
    except Exception as e:
        logger.error(f"忽略节点错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"忽略节点时出错: {str(e)}"
        }), 500

@app.teardown_appcontext
def close_db(exception):
    """关闭数据库连接"""
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()

# 添加历史记录相关API
@app.route('/api/history/save', methods=['POST'])
def save_history():
    """保存知识图谱到历史记录"""
    try:
        data = request.json
        content = data.get('content')
        title = data.get('title', f"知识图谱 {time.strftime('%Y-%m-%d %H:%M:%S')}")
        description = data.get('description', '')
        history_id = data.get('history_id')  # 如果是更新现有历史记录
        
        if not content:
            return jsonify({"success": False, "message": "内容不能为空"}), 400
        
        # 获取当前用户ID
        user_id = get_current_user_id()
        
        # 保存API响应到文件
        timestamp = time.strftime('%Y%m%d_%H%M%S')
        filename = f"api_response_{timestamp}.txt"
        
        # 确保api_responses目录存在
        api_responses_dir = 'api_responses'
        try:
            if not os.path.exists(api_responses_dir):
                os.makedirs(api_responses_dir, exist_ok=True)
                logger.info(f"创建目录: {api_responses_dir}")
        except Exception as e:
            logger.error(f"创建目录失败: {str(e)}")
            # 如果无法创建目录，使用当前目录
            api_responses_dir = '.'
        
        # 保存到api_responses目录
        file_path = os.path.join(api_responses_dir, filename)
        
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(content, f, ensure_ascii=False, indent=2)
            logger.info(f"API响应已保存到文件: {file_path}")
        except Exception as e:
            logger.error(f"保存到api_responses目录失败: {str(e)}")
            # 如果保存失败，尝试保存到当前目录
            file_path = filename
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(content, f, ensure_ascii=False, indent=2)
            logger.info(f"API响应已保存到当前目录: {file_path}")
        
        # 如果提供了历史ID，则更新记录
        if history_id:
            with get_db() as db:
                # 检查历史记录是否存在，不再按用户ID过滤
                cursor = db.execute(
                    'SELECT id FROM history_records WHERE id = ?',
                    (history_id,)
                )
                if not cursor.fetchone():
                    # 如果不存在，创建新记录
                    history_id = str(uuid.uuid4())
                    created_at = time.strftime('%Y-%m-%d %H:%M:%S')
                    db.execute(
                        'INSERT INTO history_records (id, title, content, created_at, description, file_path, user_id) VALUES (?, ?, ?, ?, ?, ?, ?)',
                        (history_id, title, json.dumps(content), created_at, description, file_path, user_id)
                    )
                else:
                    # 更新现有记录
                    db.execute(
                        'UPDATE history_records SET title = ?, content = ?, description = ?, file_path = ? WHERE id = ?',
                        (title, json.dumps(content), description, file_path, history_id)
                    )
                db.commit()
            
            return jsonify({
                "success": True, 
                "message": "历史记录更新成功", 
                "history_id": history_id,
                "file_path": file_path
            })
        else:
            # 生成唯一ID
            history_id = str(uuid.uuid4())
            created_at = time.strftime('%Y-%m-%d %H:%M:%S')
            
            # 保存到数据库
            with get_db() as db:
                db.execute(
                    'INSERT INTO history_records (id, title, content, created_at, description, file_path, user_id) VALUES (?, ?, ?, ?, ?, ?, ?)',
                    (history_id, title, json.dumps(content), created_at, description, file_path, user_id)
                )
                db.commit()
            
            # 设置cookie保存临时用户ID
            response = jsonify({
                "success": True, 
                "message": "历史记录保存成功", 
                "history_id": history_id,
                "file_path": file_path
            })
            
            # 如果是临时用户ID，设置cookie
            if not session.get('user_id'):
                response.set_cookie('temp_user_id', user_id, max_age=86400*30)  # 30天有效期
                
            return response
    
    except Exception as e:
        logger.error(f"保存历史记录失败: {str(e)}", exc_info=True)
        return jsonify({"success": False, "message": f"保存失败: {str(e)}"}), 500

@app.route('/api/history/list', methods=['GET'])
def get_history_list():
    """获取历史记录列表"""
    try:
        # 获取当前用户ID
        user_id = get_current_user_id()
        logger.info(f"获取历史记录列表，用户ID: {user_id}")
        
        # 使用应用上下文
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 获取所有历史记录，不再按用户ID过滤
            cursor.execute(
                'SELECT id, title, created_at, description FROM history_records ORDER BY created_at DESC'
            )
            
            history_records = [
                {
                    "id": row[0],
                    "title": row[1],
                    "created_at": row[2],
                    "description": row[3]
                }
                for row in cursor.fetchall()
            ]
            
            logger.info(f"找到 {len(history_records)} 条历史记录")
        
        return jsonify({
            "success": True,
            "history_records": history_records
        })
    
    except Exception as e:
        logger.error(f"获取历史记录列表失败: {str(e)}", exc_info=True)
        return jsonify({"success": False, "message": f"获取失败: {str(e)}"}), 500

@app.route('/api/history/<history_id>', methods=['GET'])
def get_history_detail(history_id):
    """获取单个历史记录详情"""
    try:
        logger.info(f"获取历史记录详情，ID: {history_id}")
        
        # 使用应用上下文
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 查询历史记录
            cursor.execute(
                'SELECT id, title, content, created_at, description, file_path FROM history_records WHERE id = ?',
                (history_id,)
            )
            row = cursor.fetchone()
            
            if not row:
                logger.error(f"历史记录不存在，ID: {history_id}")
                return jsonify({"success": False, "message": "历史记录不存在"}), 404
            
            content = None
            # 如果有文件路径，优先从文件加载内容
            file_path = row[5]  # file_path字段
            if file_path and os.path.exists(file_path):
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = json.load(f)
                    logger.info(f"从文件加载历史记录内容: {file_path}")
                except Exception as e:
                    logger.error(f"从文件加载历史记录内容失败: {str(e)}")
                    # 如果文件路径存在但无法读取，尝试在api_responses目录中查找
                    filename = os.path.basename(file_path)
                    alt_path = os.path.join('api_responses', filename)
                    if os.path.exists(alt_path):
                        try:
                            with open(alt_path, 'r', encoding='utf-8') as f:
                                content = json.load(f)
                            logger.info(f"从备用路径加载历史记录内容: {alt_path}")
                        except Exception as e2:
                            logger.error(f"从备用路径加载历史记录内容失败: {str(e2)}")
            
            # 如果从文件加载失败，则使用数据库中的内容
            if content is None and row[2]:  # content字段
                try:
                    content = json.loads(row[2])
                    logger.info("从数据库加载历史记录内容")
                except Exception as e:
                    logger.error(f"解析数据库中的历史记录内容失败: {str(e)}")
            
            if content is None:
                # 尝试在当前目录和api_responses目录中查找对应ID的文件
                possible_files = []
                for file in os.listdir('.'):
                    if file.startswith('api_response_') and file.endswith('.txt'):
                        possible_files.append(file)
                
                if os.path.exists('api_responses'):
                    for file in os.listdir('api_responses'):
                        if file.startswith('api_response_') and file.endswith('.txt'):
                            possible_files.append(os.path.join('api_responses', file))
                
                # 按修改时间排序，最新的优先
                possible_files.sort(key=lambda f: os.path.getmtime(f), reverse=True)
                
                # 尝试加载最新的文件
                if possible_files:
                    try:
                        newest_file = possible_files[0]
                        with open(newest_file, 'r', encoding='utf-8') as f:
                            content = json.load(f)
                        logger.info(f"从最新的API响应文件加载内容: {newest_file}")
                        
                        # 更新数据库中的文件路径
                        db.execute(
                            'UPDATE history_records SET file_path = ? WHERE id = ?',
                            (newest_file, history_id)
                        )
                        db.commit()
                        file_path = newest_file
                    except Exception as e:
                        logger.error(f"从最新的API响应文件加载内容失败: {str(e)}")
            
            if content is None:
                logger.error("历史记录内容为空")
                return jsonify({"success": False, "message": "历史记录内容损坏或不存在"}), 500
            
            history_record = {
                "id": row[0],
                "title": row[1],
                "content": content,
                "created_at": row[3],
                "description": row[4],
                "file_path": file_path
            }
            
            logger.info(f"成功获取历史记录详情，ID: {history_id}")
        
        return jsonify({
            "success": True,
            "history_record": history_record
        })
    
    except Exception as e:
        logger.error(f"获取历史记录详情失败: {str(e)}", exc_info=True)
        return jsonify({"success": False, "message": f"获取失败: {str(e)}"}), 500

@app.route('/api/history/<history_id>/delete', methods=['POST'])
def delete_history_record(history_id):
    """删除单个历史记录"""
    try:
        # 获取当前用户ID
        user_id = get_current_user_id()
        
        with get_db() as db:
            # 检查历史记录是否存在，不再按用户ID过滤
            cursor = db.execute(
                'SELECT id, file_path FROM history_records WHERE id = ?',
                (history_id,)
            )
            row = cursor.fetchone()
            
            if not row:
                return jsonify({"success": False, "message": "历史记录不存在"}), 404
            
            file_path = row[1]
            
            # 删除历史记录
            db.execute('DELETE FROM history_records WHERE id = ?', (history_id,))
            db.commit()
            
            # 尝试删除关联的API响应文件
            if file_path and os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    logger.info(f"删除历史记录文件: {file_path}")
                except Exception as e:
                    logger.error(f"删除历史记录文件失败: {str(e)}", exc_info=True)
            
            return jsonify({
                "success": True,
                "message": "历史记录已删除"
            })
    
    except Exception as e:
        logger.error(f"删除历史记录失败: {str(e)}", exc_info=True)
        return jsonify({"success": False, "message": f"删除失败: {str(e)}"}), 500

@app.route('/api/save_node_status', methods=['POST'])
def save_node_status():
    """保存节点点亮状态"""
    try:
        logger.info("接收到保存节点状态的请求")
        
        # 检查请求数据
        if not request.is_json:
            logger.error("请求不是JSON格式")
            return jsonify({"success": False, "message": "请求必须是JSON格式"}), 400
        
        data = request.get_json()
        logger.info(f"接收到的数据: {data.keys() if data else None}")
        
        if not data or 'graph_data' not in data:
            logger.error(f"缺少必要的数据: {data}")
            return jsonify({"success": False, "message": "缺少必要的数据"}), 400
        
        # 获取图谱数据
        graph_data = data.get('graph_data')
        node_id = data.get('node_id', 'unknown')
        highlighted = data.get('highlighted', False)
        file_title = data.get('title', f"知识图谱 {time.strftime('%Y-%m-%d %H:%M:%S')}")
        history_id = data.get('history_id', None)
        
        # 获取当前用户ID
        user_id = get_current_user_id()
        
        logger.info(f"准备保存节点 {node_id} 的点亮状态: {highlighted}")
        logger.info(f"图谱数据包含 {len(graph_data)} 条边")
        
        # 保存到文件
        timestamp = time.strftime('%Y%m%d_%H%M%S')
        filename = f"api_response_{timestamp}.txt"
        
        # 确保api_responses目录存在
        api_responses_dir = 'api_responses'
        try:
            if not os.path.exists(api_responses_dir):
                os.makedirs(api_responses_dir, exist_ok=True)
                logger.info(f"创建目录: {api_responses_dir}")
        except Exception as e:
            logger.error(f"创建目录失败: {str(e)}")
            # 如果无法创建目录，使用当前目录
            api_responses_dir = '.'
        
        # 保存到api_responses目录
        file_path = os.path.join(api_responses_dir, filename)
        
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(graph_data, f, ensure_ascii=False, indent=2)
            logger.info(f"节点状态已保存到文件: {file_path}")
        except Exception as e:
            logger.error(f"保存到api_responses目录失败: {str(e)}")
            # 如果保存失败，尝试保存到当前目录
            file_path = filename
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(graph_data, f, ensure_ascii=False, indent=2)
            logger.info(f"节点状态已保存到当前目录: {file_path}")
        
        # 更新历史记录（如果有历史ID则更新，否则创建新记录）
        try:
            with app.app_context():
                with get_db() as db:
                    if history_id:
                        # 检查历史记录是否存在，不再按用户ID过滤
                        cursor = db.execute(
                            'SELECT id FROM history_records WHERE id = ?',
                            (history_id,)
                        )
                        if cursor.fetchone():
                            # 更新现有历史记录
                            db.execute(
                                'UPDATE history_records SET content = ?, description = ?, file_path = ? WHERE id = ?',
                                (json.dumps(graph_data), f"更新于 {time.strftime('%Y-%m-%d %H:%M:%S')}", file_path, history_id)
                            )
                            logger.info(f"更新历史记录: {history_id}")
                        else:
                            # 如果历史记录不存在，创建新记录
                            history_id = str(uuid.uuid4())
                            created_at = time.strftime('%Y-%m-%d %H:%M:%S')
                            db.execute(
                                'INSERT INTO history_records (id, title, content, created_at, description, file_path, user_id) VALUES (?, ?, ?, ?, ?, ?, ?)',
                                (history_id, file_title, json.dumps(graph_data), created_at, "用户更新的知识图谱", file_path, user_id)
                            )
                            logger.info(f"创建新历史记录: {history_id}")
                    else:
                        # 创建新历史记录
                        new_history_id = str(uuid.uuid4())
                        created_at = time.strftime('%Y-%m-%d %H:%M:%S')
                        
                        db.execute(
                            'INSERT INTO history_records (id, title, content, created_at, description, file_path, user_id) VALUES (?, ?, ?, ?, ?, ?, ?)',
                            (new_history_id, file_title, json.dumps(graph_data), created_at, "用户更新的知识图谱", file_path, user_id)
                        )
                        history_id = new_history_id
                        logger.info(f"创建新历史记录: {new_history_id}")
                    
                    db.commit()
        except Exception as e:
            logger.error(f"更新历史记录失败: {str(e)}", exc_info=True)
        
        # 设置cookie保存临时用户ID
        response = jsonify({
            "success": True,
            "message": "节点状态已保存",
            "filename": filename,
            "file_path": file_path,
            "history_id": history_id
        })
        
        # 如果是临时用户ID，设置cookie
        if not session.get('user_id'):
            response.set_cookie('temp_user_id', user_id, max_age=86400*30)  # 30天有效期
        
        return response
    
    except Exception as e:
        logger.error(f"保存节点状态失败: {str(e)}", exc_info=True)
        return jsonify({"success": False, "message": f"保存失败: {str(e)}"}), 500

def get_current_user_id():
    """获取当前登录用户的ID，如果未登录则返回默认ID"""
    # 尝试从会话中获取用户ID
    user_id = session.get('user_id')
    
    # 如果未登录，使用默认ID
    if not user_id:
        # 检查是否有存储在cookie中的临时ID
        user_id = request.cookies.get('temp_user_id')
        
        # 如果没有临时ID，生成一个新的
        if not user_id:
            user_id = str(uuid.uuid4())
    
    return user_id

if __name__ == '__main__':
    # 创建必要的文件夹
    folders = ['uploads', 'static/css', 'static/js', 'templates']
    for folder in folders:
        if not os.path.exists(folder):
            os.makedirs(folder)
    
    # 初始化数据库
    try:
        init_db()
        logger.info("数据库初始化完成")
    except Exception as e:
        logger.error(f"数据库初始化异常: {str(e)}", exc_info=True)
        logger.info("尝试继续运行，但可能会出现数据库相关错误")
    
    logger.info("知识图谱生成系统启动中...")
    app.run(debug=True, port=5001)