import os
import re
import json
import uuid
import time
import threading
import logging
import requests
import sqlite3
from PyPDF2 import PdfReader
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from flask import Flask, request, jsonify, render_template, g
from flask_cors import CORS
from contextlib import closing
from collections import defaultdict

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("knowledge_graph.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("KnowledgeGraphGenerator")

# 初始化Flask应用
app = Flask(__name__, static_folder='static', template_folder='templates')
CORS(app, resources={r"/*": {"origins": "*"}})  # 增强CORS配置

# 配置上传文件夹
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# DeepSeek API配置（示例密钥，需替换为实际密钥）
OPENAI_API_KEY = "sk-ba9cc9a26b8c4859ba5c9bad33785093"
MAX_RETRIES = 3
BACKOFF_FACTOR = 2

# 数据库配置
DATABASE = os.path.join(app.root_path, 'knowledge_graph.db')

def get_db():
    """获取数据库连接"""
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
        db.row_factory = sqlite3.Row
    return db

def init_db():
    """初始化数据库（包含问答会话表）"""
    logger.info("开始初始化数据库...")
    with app.app_context():
        db_path = os.path.abspath(DATABASE)
        logger.info(f"数据库文件路径: {db_path}")
        
        # 确保数据库目录存在
        db_dir = os.path.dirname(db_path)
        if not os.path.exists(db_dir):
            os.makedirs(db_dir)
        
        # 如果数据库已存在，跳过初始化
        if os.path.exists(db_path):
            logger.info("数据库已存在，跳过初始化")
            return
        
        logger.info("数据库文件不存在，创建新数据库...")
        with closing(get_db()) as db:
            try:
                schema = """
                CREATE TABLE IF NOT EXISTS topologies (
                    id TEXT PRIMARY KEY,
                    content TEXT NOT NULL,
                    max_nodes INTEGER DEFAULT 0,
                    created_at TEXT
                );

                CREATE TABLE IF NOT EXISTS nodes (
                    id TEXT,
                    topology_id TEXT,
                    label TEXT,
                    level INTEGER,
                    value INTEGER,
                    mastered INTEGER DEFAULT 0,
                    mastery_score REAL DEFAULT 0,
                    consecutive_correct INTEGER DEFAULT 0,
                    content_snippet TEXT,
                    PRIMARY KEY (topology_id, id),
                    FOREIGN KEY (topology_id) REFERENCES topologies (id)
                );

                CREATE TABLE IF NOT EXISTS edges (
                    topology_id TEXT,
                    from_node TEXT,
                    to_node TEXT,
                    label TEXT,
                    PRIMARY KEY (topology_id, from_node, to_node),
                    FOREIGN KEY (topology_id) REFERENCES topologies (id),
                    FOREIGN KEY (from_node) REFERENCES nodes (id),
                    FOREIGN KEY (to_node) REFERENCES nodes (id)
                );

                CREATE TABLE IF NOT EXISTS questions (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    node_id TEXT,
                    question TEXT,
                    answer TEXT,
                    feedback TEXT,
                    correctness INTEGER DEFAULT 0,
                    created_at TEXT,
                    answered_at TEXT,
                    session_id TEXT,
                    FOREIGN KEY (topology_id, node_id) REFERENCES nodes (topology_id, id)
                );

                -- 问答会话表
                CREATE TABLE IF NOT EXISTS quiz_sessions (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    node_id TEXT,
                    created_at TEXT,
                    questions_answered INTEGER DEFAULT 0,
                    consecutive_correct INTEGER DEFAULT 0,
                    mastered INTEGER DEFAULT 0,
                    FOREIGN KEY (topology_id, node_id) REFERENCES nodes (topology_id, id)
                );
                """
                db.cursor().executescript(schema)
                db.commit()
                logger.info("数据库表创建成功")
            except Exception as e:
                logger.error(f"数据库初始化失败: {str(e)}", exc_info=True)
                raise

def clean_json_string(s: str) -> str:
    """清洗模型输出，去除Markdown代码块标记"""
    s = re.sub(r"```(?:json)?", "", s)
    return s.strip()

def enhance_json_format(json_str: str) -> str:
    """增强JSON格式，处理各种复杂格式问题"""
    import json
    import re
    
    logger.info(f"原始JSON内容: {json_str[:200]}...")
    
    # 1. 移除可能的前导/尾随非JSON字符
    json_str = re.sub(r'^.*?(\[.*\]).*$', r'\1', json_str, flags=re.DOTALL)
    
    # 2. 处理行首缩进和换行
    json_str = re.sub(r'\n\s*', ' ', json_str)
    
    # 3. 处理可能的单引号问题
    json_str = re.sub(r"'", '"', json_str)
    
    # 4. 处理常见的格式错误（尝试修复）
    try:
        # 尝试直接解析
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        logger.error(f"初始解析失败: {e}")
        
        # 错误定位
        error_pos = e.pos
        logger.error(f"错误位置: {error_pos}, 附近内容: {json_str[error_pos-20:error_pos+20]}")
        
        # 5. 智能修复常见错误
        # 5.1 处理缺少逗号的情况
        if e.msg == "Expecting ',' delimiter":
            logger.info("尝试修复缺少逗号的问题...")
            json_list = list(json_str)
            
            # 在错误位置前查找可能缺少逗号的位置
            search_start = max(0, error_pos - 100)
            bracket_count = 0
            for i in range(error_pos-1, search_start, -1):
                if json_list[i] == ']':
                    bracket_count += 1
                elif json_list[i] == '[':
                    bracket_count -= 1
                
                # 找到合适的位置插入逗号
                if bracket_count == 0 and json_list[i] == ']':
                    json_list.insert(i+1, ',')
                    logger.info(f"在位置 {i+1} 插入逗号")
                    break
            
            # 尝试再次解析
            try:
                return json.loads(''.join(json_list))
            except json.JSONDecodeError as e2:
                logger.error(f"修复后仍失败: {e2}")
        
        # 5.2 处理未闭合的引号
        if e.msg.startswith('Expecting property name enclosed in double quotes'):
            logger.info("尝试修复未闭合的引号...")
            # 查找最近的未闭合引号并添加
            unclosed_quote_pos = json_str.rfind('"', 0, error_pos)
            if unclosed_quote_pos != -1:
                json_list = list(json_str)
                json_list.insert(error_pos, '"')
                logger.info(f"在位置 {error_pos} 添加引号")
                try:
                    return json.loads(''.join(json_list))
                except json.JSONDecodeError as e3:
                    logger.error(f"修复后仍失败: {e3}")
        
        # 6. 使用更宽松的解析器（作为最后的手段）
        try:
            import ast
            logger.info("尝试使用ast.literal_eval进行宽松解析...")
            parsed = ast.literal_eval(json_str)
            # 转换为标准JSON
            return json.loads(json.dumps(parsed))
        except Exception as e4:
            logger.error(f"宽松解析失败: {e4}")
            logger.error(f"无法修复的JSON内容: {json_str}")
            raise RuntimeError("无法解析API返回的JSON格式") from e4

def sanitize_text(text: str) -> str:
    """清理文本中的特殊字符，防止破坏JSON格式"""
    # 移除可能干扰JSON解析的字符
    text = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', text)  # 移除控制字符
    # 转义特殊字符
    text = text.replace('\\', '\\\\')
    text = text.replace('"', '\\"')
    return text

def extract_knowledge_from_text(text: str, max_nodes: int = 0, max_retries: int = MAX_RETRIES) -> list:
    """调用DeepSeek API提取适合树形结构的知识点层级关系"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 清理文本
    sanitized_text = sanitize_text(text)
    
    # 根据节点数量限制调整提示
    node_limit_prompt = ""
    if max_nodes > 0:
        node_limit_prompt = f"请确保最终提取的知识点数量不超过{max_nodes}个。"
    
    messages = [
        {"role": "system", "content": f"""你是一个知识图谱构建专家，能够从文本中提取知识点并构建树形结构。
请分析文本内容，识别出主要知识点及其层级关系（如父节点-子节点关系），
以JSON数组形式输出，每个元素格式为 [父知识点, 关系, 子知识点]。
关系应体现层级结构，如"包含"、"属于"、"是子类"等。确保输出格式正确，仅返回JSON数组。
{node_limit_prompt}"""},
        {"role": "user", "content": f"请从下面文本中提取知识点及其层级关系，输出JSON数组，每个元素格式为 [父知识点, 关系, 子知识点]：\n{sanitized_text}"}
    ]
    
    backoff = 2
    for attempt in range(1, max_retries + 1):
        try:
            logger.info(f"第{attempt}次尝试调用DeepSeek API...")
            
            # 使用OpenAI SDK调用API
            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=messages,
                max_tokens=1500
            )
            
            raw = response.choices[0].message.content
            cleaned = clean_json_string(raw)
            logger.info(f"API返回知识关系: {cleaned[:200]}...")
            
            # 保存原始响应用于调试
            with open(f"api_response_{time.time()}.txt", "w", encoding="utf-8") as f:
                f.write(cleaned)
            
            # 增强JSON解析
            knowledge_edges = enhance_json_format(cleaned)
            
            # 验证输出格式
            if not isinstance(knowledge_edges, list):
                raise ValueError(f"API返回非数组格式: {type(knowledge_edges)}")
            for idx, item in enumerate(knowledge_edges):
                if not (isinstance(item, list) and len(item) == 3):
                    raise ValueError(f"API返回元素格式错误，应为三元组，位置 {idx}: {item}")
                    
            logger.info(f"成功解析知识关系，共{len(knowledge_edges)}条")
            return knowledge_edges
                
        except Exception as e:
            logger.error(f"API请求错误: {str(e)}", exc_info=True)
            if attempt < max_retries:
                logger.info(f"准备第{attempt+1}次重试...")
                time.sleep(backoff)
                backoff *= BACKOFF_FACTOR
            else:
                raise
    
    raise RuntimeError("多次重试后仍无法获取有效响应")

def parse_document(file_path):
    """解析文档内容，返回文本（新增PPT支持）"""
    file_ext = os.path.splitext(file_path)[1].lower()
    logger.info(f"开始解析文档: {file_path}, 类型: {file_ext}")
    
    try:
        if file_ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        elif file_ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    text += page_text
                    if page_num % 10 == 0:
                        logger.info(f"已解析PDF第 {page_num} 页")
            return text
        elif file_ext in ['.docx', '.doc']:
            doc = Document(file_path)
            full_text = []
            for para_num, para in enumerate(doc.paragraphs):
                full_text.append(para.text)
                if para_num % 50 == 0:
                    logger.info(f"已解析Word第 {para_num} 段落")
            return '\n'.join(full_text)
        elif file_ext == '.html':
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            soup = BeautifulSoup(html_content, 'lxml')
            text = soup.get_text()
            return ' '.join(text.split())
        elif file_ext in ['.ppt', '.pptx']:  # 新增PPT/PPTX支持
            text = ""
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if slide_num % 10 == 0:
                    logger.info(f"已解析PPT第 {slide_num} 页")
            return text
        else:
            logger.error(f"不支持的文件格式: {file_ext}")
            return None
    except Exception as e:
        logger.error(f"解析文档出错: {str(e)}", exc_info=True)
        return None

def extract_content_snippet(content: str, topic: str) -> str:
    """从原文中提取与主题相关的片段"""
    # 查找主题在内容中的位置
    index = content.lower().find(topic.lower())
    if index == -1:
        return ""
    
    # 提取上下文片段（主题前后各200个字符）
    start = max(0, index - 200)
    end = min(len(content), index + len(topic) + 200)
    snippet = content[start:end]
    
    # 确保片段包含主题
    if topic.lower() not in snippet.lower():
        return ""
    
    # 添加省略号表示截断
    if start > 0:
        snippet = "..." + snippet
    if end < len(content):
        snippet = snippet + "..."
    
    return snippet

def build_tree_structure(knowledge_edges, topology_id, content: str, max_nodes: int = 0):
    """构建树形知识图数据结构，保存原文片段并恢复掌握状态"""
    nodes = {}
    edges = []
    all_node_ids = set()
    
    # 收集所有节点
    for src, rel, tgt in knowledge_edges:
        all_node_ids.add(src)
        all_node_ids.add(tgt)
        
        # 确保节点存在
        if src not in nodes:
            # 提取原文片段
            snippet = extract_content_snippet(content, src)
            nodes[src] = {
                "id": src,
                "label": src,
                "title": src,  # 鼠标悬停时显示的完整标题
                "level": 0,    # 默认层级
                "value": 1,    # 节点大小基准
                "mastered": False,  # 知识点掌握状态
                "mastery_score": 0,  # 掌握分数
                "consecutive_correct": 0,  # 连续正确回答次数
                "content_snippet": snippet  # 保存原文片段
            }
        if tgt not in nodes:
            snippet = extract_content_snippet(content, tgt)
            nodes[tgt] = {
                "id": tgt,
                "label": tgt,
                "title": tgt,
                "level": 0,
                "value": 1,
                "mastered": False,
                "mastery_score": 0,
                "consecutive_correct": 0,
                "content_snippet": snippet
            }
        
        # 添加边
        edges.append({
            "from": src,
            "to": tgt,
            "label": rel,
            "title": rel,  # 鼠标悬停时显示的边关系
            "arrows": "to",
            "font": {
                "align": "middle"
            }
        })
    
    # 计算节点层级
    def calculate_level(node_id, current_level=0):
        if node_id in nodes:
            nodes[node_id]["level"] = max(nodes[node_id]["level"], current_level)
            # 递归设置子节点层级
            for edge in edges:
                if edge["from"] == node_id:
                    calculate_level(edge["to"], current_level + 1)
    
    # 找到根节点（没有父节点的节点）
    root_candidates = all_node_ids.copy()
    for _, _, tgt in knowledge_edges:
        root_candidates.discard(tgt)
    
    root = next(iter(root_candidates)) if root_candidates else (list(nodes.keys())[0] if nodes else None)
    
    # 从根节点开始计算层级
    if root:
        calculate_level(root)
    
    # 计算节点重要性（连接数）
    for node_id in nodes:
        in_connections = sum(1 for edge in edges if edge["to"] == node_id)
        out_connections = sum(1 for edge in edges if edge["from"] == node_id)
        nodes[node_id]["value"] = max(1, in_connections + out_connections)
        
        # 从数据库获取并恢复节点的掌握状态
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT mastered, mastery_score, consecutive_correct FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node_status = cursor.fetchone()
            if node_status:
                nodes[node_id]["mastered"] = bool(node_status["mastered"])
                nodes[node_id]["mastery_score"] = node_status["mastery_score"]
                nodes[node_id]["consecutive_correct"] = node_status["consecutive_correct"]
    
    # 保存节点和边到数据库
    save_to_database(topology_id, list(nodes.values()), edges, content, max_nodes)
    
    # 转换为节点列表
    tree_nodes = list(nodes.values())
    
    return {
        "nodes": tree_nodes,
        "edges": edges,
        "root": root
    }

def save_to_database(topology_id, nodes, edges, content: str, max_nodes=0):
    """将知识图谱数据保存到数据库（保存原文内容和节点数量限制）"""
    with app.app_context():
        db = get_db()
        cursor = db.cursor()
        
        # 保存拓扑图信息（包含原文内容和节点数量限制）
        cursor.execute(
            "INSERT OR REPLACE INTO topologies (id, content, max_nodes, created_at) VALUES (?, ?, ?, ?)",
            (topology_id, content, max_nodes, time.strftime('%Y-%m-%d %H:%M:%S'))
        )
        
        # 保存节点
        for node in nodes:
            cursor.execute(
                """INSERT OR REPLACE INTO nodes 
                (topology_id, id, label, level, value, mastered, mastery_score, consecutive_correct, content_snippet) 
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (topology_id, node["id"], node["label"], node["level"], node["value"], 
                 int(node["mastered"]), node["mastery_score"], node["consecutive_correct"], 
                 node.get("content_snippet", ""))
            )
        
        # 保存边
        for edge in edges:
            cursor.execute(
                "INSERT OR REPLACE INTO edges (topology_id, from_node, to_node, label) VALUES (?, ?, ?, ?)",
                (topology_id, edge["from"], edge["to"], edge["label"])
            )
        
        db.commit()

def process_document(file_path, topology_id, max_nodes=0):
    """处理文档并生成树形知识图（支持节点数量限制）"""
    start_time = time.time()
    logger.info(f"开始处理文档: {file_path}, 拓扑ID: {topology_id}, 最大节点数: {max_nodes}")
    
    # 更新状态为处理中
    topology_results[topology_id] = {
        "status": "processing",
        "progress": 0,
        "message": "开始处理文档...",
        "created_at": time.strftime('%Y-%m-%d %H:%M:%S')
    }
    
    try:
        with app.app_context():
            update_progress(topology_id, 10, "解析文档内容...")
            text = parse_document(file_path)
            if not text:
                topology_results[topology_id] = {
                    "status": "error",
                    "message": "无法解析文档内容"
                }
                logger.error(f"文档解析失败: {file_path}")
                return
            update_progress(topology_id, 20, "准备提取知识层级...")
            text_length = len(text)
            if text_length < 100:
                topology_results[topology_id] = {
                    "status": "error",
                    "message": "文档内容过短，无法提取知识"
                }
                logger.warning(f"文档内容过短: {file_path}, 长度: {text_length}")
                return

            MAX_TEXT_LENGTH = 8000
            if text_length > MAX_TEXT_LENGTH:
                logger.info(f"文档内容过长，截取前{MAX_TEXT_LENGTH}字符")
                text = text[:MAX_TEXT_LENGTH]
            
            update_progress(topology_id, 60, "调用DeepSeek API提取知识层级...")
            knowledge_edges = extract_knowledge_from_text(text, max_nodes)
            logger.info(f"成功提取{len(knowledge_edges)}条知识层级关系")
            
            update_progress(topology_id, 80, "构建树形知识图并提取原文片段...")
            knowledge_graph = build_tree_structure(knowledge_edges, topology_id, text, max_nodes)
            
            processing_time = time.time() - start_time
            logger.info(f"树形知识图生成完成，耗时: {processing_time:.2f} 秒")
            
            update_progress(topology_id, 100, "处理完成")
            time.sleep(1)  # 确保前端有足够时间更新进度
            
            topology_results[topology_id] = {
                "status": "completed",
                "data": knowledge_graph,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph["nodes"]),
                "edge_count": len(knowledge_graph["edges"]),
                "processing_time": round(processing_time, 2),
                "text_length": text_length,
                "max_nodes": max_nodes  # 保存节点数量限制
            }
            
    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        logger.error(f"处理文档出错: {str(e)}", exc_info=True)
        topology_results[topology_id] = {
            "status": "error",
            "message": f"处理过程中出错: {str(e)}"
        }

def update_progress(topology_id, progress, message):
    """更新处理进度"""
    if topology_id in topology_results:
        topology_results[topology_id].update({
            'progress': progress,
            'message': message
        })
        logger.info(f"拓扑ID: {topology_id}, 进度: {progress}%, 消息: {message}")

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/generate', methods=['POST'])
def generate_knowledge_graph():
    """处理用户点击生成按钮后的请求，支持节点数量控制"""
    if 'file' not in request.files:
        logger.error("文件上传错误: 没有文件")
        return jsonify({'status': 'error', 'message': '没有文件上传'}), 400
    
    file = request.files['file']
    if file.filename == '':
        logger.error("文件上传错误: 未选择文件")
        return jsonify({'status': 'error', 'message': '未选择文件'}), 400
    
    # 获取节点数量 - 确保从表单获取
    max_nodes = request.form.get('max_nodes', 0, type=int)
    
    # 检查文件大小
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    
    if file_size > 50 * 1024 * 1024:  # 50MB限制
        logger.error(f"文件上传错误: 文件过大 ({file_size/1024/1024:.2f} MB)")
        return jsonify({'status': 'error', 'message': '文件大小超过50MB限制'}), 400
    
    topology_id = str(uuid.uuid4())
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{topology_id}_{file.filename}")
    file.save(file_path)
    
    logger.info(f"文件上传成功: {file_path}, 大小: {file_size/1024/1024:.2f} MB, 最大节点数: {max_nodes}")
    
    # 启动处理线程，并在应用上下文中执行
    threading.Thread(
        target=lambda: with_app_context(process_document, file_path, topology_id, max_nodes)
    ).start()
    
    return jsonify({
        'status': 'success',
        'topology_id': topology_id,
        'message': '文档上传成功，正在生成知识图谱',
        'max_nodes': max_nodes  # 返回节点数量限制
    })

def with_app_context(func, *args, **kwargs):
    """在应用上下文中执行函数"""
    with app.app_context():
        func(*args, **kwargs)

@app.route('/api/topology/<topology_id>', methods=['GET'])
def get_topology(topology_id):
    if topology_id not in topology_results:
        # 尝试从数据库获取
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT id, content, max_nodes, created_at FROM topologies WHERE id = ?",
                (topology_id,)
            )
            topology = cursor.fetchone()
            
            if not topology:
                logger.error(f"获取拓扑图错误: ID不存在 ({topology_id})")
                return jsonify({'status': 'error', 'message': '拓扑图不存在'}), 404
            
            # 从数据库获取节点和边
            cursor.execute(
                "SELECT id, label, level, value, mastered, mastery_score, consecutive_correct, content_snippet FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            nodes = [dict(row) for row in cursor.fetchall()]
            
            cursor.execute(
                "SELECT from_node, to_node, label FROM edges WHERE topology_id = ?",
                (topology_id,)
            )
            edges = [dict(row) for row in cursor.fetchall()]
            
            knowledge_graph = {
                "nodes": nodes,
                "edges": edges,
                "root": next((node["id"] for node in nodes if node["level"] == 0), nodes[0]["id"] if nodes else None)
            }
            
            return jsonify({
                'status': 'success',
                'data': knowledge_graph,
                'created_at': topology["created_at"],
                'node_count': len(nodes),
                'edge_count': len(edges),
                'text_length': len(topology["content"]),
                'max_nodes': topology["max_nodes"]  # 返回节点数量限制
            })
    
    topology = topology_results[topology_id]
    
    if topology['status'] == 'processing':
        return jsonify({
            'status': 'processing',
            'progress': topology.get('progress', 0),
            'message': topology.get('message', '正在处理中'),
            'max_nodes': topology.get('max_nodes', 0)  # 返回节点数量限制
        })
    
    if topology['status'] == 'error':
        logger.error(f"获取拓扑图错误: {topology.get('message', '未知错误')}")
        return jsonify({
            'status': 'error',
            'message': topology.get('message', '生成知识图时出错')
        }), 500
    
    return jsonify({
        'status': 'success',
        'data': topology['data'],
        'created_at': topology['created_at'],
        'node_count': topology['node_count'],
        'edge_count': topology['edge_count'],
        'processing_time': topology['processing_time'],
        'text_length': topology.get('text_length', 0),
        'max_nodes': topology.get('max_nodes', 0)  # 返回节点数量限制
    })

@app.route('/api/topology/<topology_id>/set_max_nodes', methods=['POST'])
def set_topology_max_nodes(topology_id):
    """更新拓扑图的节点数量设置"""
    try:
        data = request.json
        max_nodes = data.get('max_nodes', 0)
        
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 更新拓扑图的节点数量设置
            cursor.execute(
                "UPDATE topologies SET max_nodes = ? WHERE id = ?",
                (max_nodes, topology_id)
            )
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'message': '节点数量设置已更新',
                'max_nodes': max_nodes
            })
            
    except Exception as e:
        logger.error(f"设置节点数量错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"设置节点数量时出错: {str(e)}"
        }), 500

@app.route('/api/topology/<topology_id>/regenerate', methods=['POST'])
def regenerate_topology(topology_id):
    """重新生成知识图谱，使用用户输入的新节点数量"""
    try:
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 从请求中获取新的节点数量
            data = request.json
            max_nodes = data.get('max_nodes', 0)  # 从请求中获取新的节点数量
            
            # 从数据库获取原文内容
            cursor.execute(
                "SELECT content FROM topologies WHERE id = ?",
                (topology_id,)
            )
            topology = cursor.fetchone()
            
            if not topology:
                return jsonify({
                    'status': 'error',
                    'message': '拓扑图不存在'
                }), 404
            
            content = topology["content"]
            
            update_progress(topology_id, 30, "重新提取知识层级...")
            knowledge_edges = extract_knowledge_from_text(content, max_nodes)  # 使用新的节点数量
            logger.info(f"重新生成成功提取{len(knowledge_edges)}条知识层级关系")
            
            # 保存当前节点的掌握状态
            cursor.execute(
                "SELECT id, mastered, mastery_score, consecutive_correct FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            mastery_states = {row["id"]: dict(row) for row in cursor.fetchall()}
            
            update_progress(topology_id, 70, "重新构建树形知识图...")
            knowledge_graph = build_tree_structure(knowledge_edges, topology_id, content, max_nodes)  # 使用新的节点数量
            
            # 恢复节点的掌握状态
            for node in knowledge_graph["nodes"]:
                node_id = node["id"]
                if node_id in mastery_states:
                    state = mastery_states[node_id]
                    node["mastered"] = bool(state["mastered"])
                    node["mastery_score"] = state["mastery_score"]
                    node["consecutive_correct"] = state["consecutive_correct"]
                    
                    # 更新数据库中的节点状态
                    cursor.execute(
                        """UPDATE nodes SET 
                        mastered = ?, mastery_score = ?, consecutive_correct = ?
                        WHERE topology_id = ? AND id = ?""",
                        (int(state["mastered"]), state["mastery_score"], 
                         state["consecutive_correct"], topology_id, node_id)
                    )
            
            # 更新拓扑图的节点数量设置到数据库
            cursor.execute(
                "UPDATE topologies SET max_nodes = ? WHERE id = ?",
                (max_nodes, topology_id)
            )
            
            # 更新处理结果
            topology_results[topology_id] = {
                "status": "completed",
                "data": knowledge_graph,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph["nodes"]),
                "edge_count": len(knowledge_graph["edges"]),
                "processing_time": 0,
                "text_length": len(content),
                "max_nodes": max_nodes  # 保存新的节点数量限制
            }
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'message': '知识图谱重新生成成功',
                'node_count': len(knowledge_graph["nodes"]),
                'edge_count': len(knowledge_graph["edges"]),
                'max_nodes': max_nodes  # 返回新的节点数量限制
            })
            
    except Exception as e:
        logger.error(f"重新生成知识图谱错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"重新生成知识图谱时出错: {str(e)}"
        }), 500


@app.route('/api/topology/<topology_id>/node/<node_id>/question', methods=['GET'])
def get_question(topology_id, node_id):
    """获取关于指定节点的问题（基于原文内容，支持会话管理）"""
    try:
        with app.app_context():
            # 从数据库获取节点信息
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            
            if not node:
                return jsonify({
                    'status': 'error',
                    'message': '节点不存在'
                }), 404
            
            node_label = node["label"]
            content_snippet = node["content_snippet"]
            
            # 检查是否已有活跃会话
            session_id = request.args.get('session_id')
            if session_id:
                cursor.execute(
                    "SELECT mastered, consecutive_correct FROM quiz_sessions WHERE id = ?",
                    (session_id,)
                )
                session = cursor.fetchone()
                if session and session["mastered"]:
                    return jsonify({
                        'status': 'success',
                        'mastered': True,
                        'message': '该知识点已掌握'
                    })
            
            # 创建或获取会话
            if not session_id:
                session_id = str(uuid.uuid4())
                cursor.execute(
                    "INSERT INTO quiz_sessions (id, topology_id, node_id, created_at) VALUES (?, ?, ?, ?)",
                    (session_id, topology_id, node_id, time.strftime('%Y-%m-%d %H:%M:%S'))
                )
                db.commit()
            
            # 获取会话状态
            cursor.execute(
                "SELECT consecutive_correct FROM quiz_sessions WHERE id = ?",
                (session_id,)
            )
            session = cursor.fetchone()
            consecutive_correct = session["consecutive_correct"] if session else 0
            
            # 生成问题（基于会话状态）
            question = generate_question(node_label, content_snippet, consecutive_correct)
            
            # 保存问题到数据库
            question_id = str(uuid.uuid4())
            cursor.execute(
                "INSERT INTO questions (id, topology_id, node_id, question, session_id, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                (question_id, topology_id, node_id, question, session_id, time.strftime('%Y-%m-%d %H:%M:%S'))
            )
            db.commit()
            
            return jsonify({
                'status': 'success',
                'data': {
                    'question_id': question_id,
                    'question': question,
                    'node_id': node_id,
                    'session_id': session_id
                }
            })
    except Exception as e:
        logger.error(f"获取问题错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"生成问题时出错: {str(e)}"
        }), 500

def generate_question(topic, context, consecutive_correct=0):
    """根据连续正确次数生成不同难度的问题"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 根据掌握程度生成不同难度的问题
    difficulty_map = {
        0: "基础概念题，用简洁的语言解释",
        1: "理解应用题，结合实例说明",
        2: "综合分析题，比较相关概念"
    }
    difficulty = difficulty_map.get(consecutive_correct, "进阶思考题，拓展应用场景")
    
    # 构建提示词
    messages = [
        {"role": "system", "content": f"""你是一个教育专家，能够基于原文内容生成有针对性的问题。
请生成一个{difficulty}，测试用户对"{topic}"的理解。
问题应该清晰明确，基于提供的原文内容，难度适合当前掌握程度。
只返回问题文本，不需要其他内容。"""},
        {"role": "user", "content": f"原文片段: {context}\n问题:"}
    ]
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=150
        )
        
        question = response.choices[0].message.content.strip()
        return question
    except Exception as e:
        logger.error(f"生成问题出错: {str(e)}", exc_info=True)
        return f"关于{topic}的问题（基于原文）"

@app.route('/api/topology/<topology_id>/question/<question_id>/answer', methods=['POST'])
def answer_question(topology_id, question_id):
    """处理用户对问题的回答（支持会话状态管理）并更新节点状态"""
    try:
        with app.app_context():
            data = request.json
            answer = data.get('answer', '')
            node_id = data.get('node_id', '')
            session_id = data.get('session_id', '')
            
            if not answer or not node_id or not session_id:
                return jsonify({
                    'status': 'error',
                    'message': '缺少必要的参数'
                }), 400
            
            # 从数据库获取问题
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT question, node_id, session_id FROM questions WHERE id = ? AND topology_id = ?",
                (question_id, topology_id)
            )
            question_data = cursor.fetchone()
            
            if not question_data:
                return jsonify({
                    'status': 'error',
                    'message': '问题不存在'
                }), 404
            
            question = question_data["question"]
            stored_node_id = question_data["node_id"]
            stored_session_id = question_data["session_id"]
            
            if stored_node_id != node_id or stored_session_id != session_id:
                return jsonify({
                    'status': 'error',
                    'message': '问题与会话不匹配'
                }), 400
            
            # 从数据库获取节点信息
            cursor.execute(
                "SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            if not node:
                return jsonify({
                    'status': 'error',
                    'message': '节点不存在'
                }), 400
            
            node_label = node["label"]
            content_snippet = node["content_snippet"]
            
            # 调用DeepSeek评估回答
            evaluation = evaluate_answer(question, answer, node_label, content_snippet)
            
            # 确定回答是否正确
            is_correct = evaluation["correct"] if "correct" in evaluation else False
            feedback_text = evaluation["feedback"] if "feedback" in evaluation else "无法评估回答"
            
            # 更新问题状态
            cursor.execute(
                "UPDATE questions SET answered_at = ?, answer = ?, feedback = ?, correctness = ? WHERE id = ?",
                (time.strftime('%Y-%m-%d %H:%M:%S'), answer, feedback_text, 1 if is_correct else 0, question_id)
            )
            
            # 更新会话状态
            cursor.execute(
                "SELECT consecutive_correct, mastered FROM quiz_sessions WHERE id = ?",
                (session_id,)
            )
            session = cursor.fetchone()
            if not session:
                return jsonify({
                    'status': 'error',
                    'message': '问答会话不存在'
                }), 404
            
            current_consecutive = session["consecutive_correct"]
            current_mastered = session["mastered"]
            
            # 更新连续正确计数
            new_consecutive = current_consecutive + 1 if is_correct else 0
            new_mastered = 1 if new_consecutive >= 3 else 0
            
            cursor.execute(
                """UPDATE quiz_sessions SET 
                questions_answered = questions_answered + 1,
                consecutive_correct = ?,
                mastered = ?
                WHERE id = ?""",
                (new_consecutive, new_mastered, session_id)
            )
            
            # 更新节点的掌握状态
            cursor.execute(
                "SELECT mastery_score, consecutive_correct FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node_status = cursor.fetchone()
            current_node_score = node_status["mastery_score"] if node_status else 0
            current_node_consecutive = node_status["consecutive_correct"] if node_status else 0
            
            node_new_score = min(10, current_node_score + (1 if is_correct else -0.5))
            node_new_consecutive = new_consecutive
            node_new_mastered = new_mastered
            
            cursor.execute(
                "UPDATE nodes SET mastery_score = ?, consecutive_correct = ?, mastered = ? WHERE topology_id = ? AND id = ?",
                (node_new_score, node_new_consecutive, int(node_new_mastered), topology_id, node_id)
            )
            
            # 确保掌握状态正确更新（显式处理）
            if node_new_mastered:
                cursor.execute(
                    "UPDATE nodes SET mastered = 1 WHERE topology_id = ? AND id = ?",
                    (topology_id, node_id)
                )
            
            # 如果未掌握，生成下一个问题
            next_question = None
            next_question_id = None
            if not new_mastered:
                # 生成下一个问题（基于更新后的状态）
                next_question = generate_question(node_label, content_snippet, new_consecutive)
                if next_question:
                    next_question_id = str(uuid.uuid4())
                    cursor.execute(
                        "INSERT INTO questions (id, topology_id, node_id, question, session_id, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                        (next_question_id, topology_id, node_id, next_question, session_id, time.strftime('%Y-%m-%d %H:%M:%S'))
                    )
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'data': {
                    'correct': is_correct,
                    'feedback': feedback_text,
                    'mastered': new_mastered,
                    'consecutive_correct': new_consecutive,
                    'session_id': session_id,
                    'next_question': {
                        'id': next_question_id,
                        'question': next_question
                    } if next_question else None
                }
            })
    except Exception as e:
        logger.error(f"处理回答错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"处理回答时出错: {str(e)}"
        }), 500

def evaluate_answer(question, answer, topic, context):
    """调用DeepSeek评估回答是否正确（包含上下文）"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    messages = [
        {"role": "system", "content": """你是一个知识评估专家，能够准确判断用户回答的正确性，并给出详细反馈。
请评估用户对问题"{question}"的回答"{answer}"是否正确，参考原文片段："{context}"。
你的评估应该包括：
1. 判断回答是否正确（是/否）
2. 提供具体的反馈，解释为什么正确或错误
3. 如果回答错误，提供正确的信息
4. 如果回答正确，考虑是否需要进一步提问以深化理解

输出格式为JSON，包含以下字段：
{
  "correct": true/false,
  "feedback": "具体的反馈信息",
  "next_question": "如果需要进一步提问，这里是下一个问题，否则为null"
}""".replace("{question}", question).replace("{answer}", answer).replace("{context}", context)},
        {"role": "user", "content": f"问题: {question}\n回答: {answer}\n原文片段: {context}\n请评估这个回答是否正确，并提供反馈。如果正确，考虑是否需要进一步提问。"}
    ]
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=300
        )
        
        response_text = response.choices[0].message.content.strip()
        logger.info(f"评估回答响应: {response_text[:200]}...")
        
        # 解析响应为JSON
        try:
            feedback = json.loads(clean_json_string(response_text))
            # 确保包含所有必要字段
            if 'correct' not in feedback:
                feedback['correct'] = False
            if 'feedback' not in feedback:
                feedback['feedback'] = "无法解析评估结果，请重试"
            if 'next_question' not in feedback:
                feedback['next_question'] = None
        except json.JSONDecodeError:
            # 如果解析失败，创建默认反馈
            feedback = {
                "correct": False,
                "feedback": "无法解析评估结果，请重试",
                "next_question": None
            }
        
        return feedback
    except Exception as e:
        logger.error(f"评估回答出错: {str(e)}", exc_info=True)
        return {
            "correct": False,
            "feedback": f"评估回答时出错: {str(e)}",
            "next_question": None
        }

@app.route('/api/topology/<topology_id>/ignore', methods=['POST'])
def ignore_nodes(topology_id):
    """忽略用户选择的节点"""
    try:
        with app.app_context():
            data = request.json
            ignored_nodes = data.get('ignored_nodes', [])
            
            db = get_db()
            cursor = db.cursor()
            
            # 获取所有节点
            cursor.execute(
                "SELECT id, label, level, value, mastered, mastery_score, content_snippet FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            all_nodes = [dict(row) for row in cursor.fetchall()]
            
            # 获取所有边
            cursor.execute(
                "SELECT from_node, to_node, label FROM edges WHERE topology_id = ?",
                (topology_id,)
            )
            all_edges = [dict(row) for row in cursor.fetchall()]
            
            # 筛选节点（排除被忽略的）
            filtered_nodes = [node for node in all_nodes if node["id"] not in ignored_nodes]
            
            # 筛选边（只保留未被忽略节点之间的边）
            filtered_edges = []
            node_ids = set([node["id"] for node in filtered_nodes])
            for edge in all_edges:
                if edge["from_node"] in node_ids and edge["to_node"] in node_ids:
                    filtered_edges.append(edge)
            
            return jsonify({
                'status': 'success',
                'data': {
                    'nodes': filtered_nodes,
                    'edges': filtered_edges,
                    'root': next((node["id"] for node in filtered_nodes if node["level"] == 0), filtered_nodes[0]["id"] if filtered_nodes else None)
                }
            })
    except Exception as e:
        logger.error(f"忽略节点错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"忽略节点时出错: {str(e)}"
        }), 500

@app.teardown_appcontext
def close_db(exception):
    """关闭数据库连接"""
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()

if __name__ == '__main__':
    # 创建必要的文件夹
    folders = ['uploads', 'static/css', 'static/js', 'templates']
    for folder in folders:
        if not os.path.exists(folder):
            os.makedirs(folder)
    
    # 初始化数据库
    try:
        init_db()
        logger.info("数据库初始化完成")
    except Exception as e:
        logger.error(f"数据库初始化异常: {str(e)}", exc_info=True)
        logger.info("尝试继续运行，但可能会出现数据库相关错误")
    
    # 拓扑图处理结果存储（使用全局变量）
    topology_results = {}
    
    logger.info("知识图谱生成系统启动中...")
    app.run(debug=True, port=5001)