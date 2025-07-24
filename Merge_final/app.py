import os
import re
import json
import uuid
import time
import threading
import logging
import sqlite3
import random
import string
import secrets
from datetime import datetime, timedelta
from functools import wraps
from PyPDF2 import PdfReader
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from flask import Flask, request, jsonify, render_template, g, session, redirect, url_for, flash
from flask_cors import CORS
from contextlib import closing
from collections import defaultdict
from flask_mail import Mail, Message
from werkzeug.security import generate_password_hash, check_password_hash

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("knowledge_graph.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("KnowledgeGraphGenerator")

# 初始化Flask应用
app = Flask(__name__, static_folder='static', template_folder='templates')
CORS(app, resources={r"/*": {"origins": "*"}})  # 增强CORS配置

# 配置上传文件夹
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# DeepSeek API配置（示例密钥，需替换为实际密钥）
OPENAI_API_KEY = "sk-420816a747e6490d8980d0807b9b5b24"
MAX_RETRIES = 3
BACKOFF_FACTOR = 2

# 数据库配置
DATABASE = os.path.join(app.root_path, 'knowledge_graph.db')

# 邮件配置 - QQ邮箱SMTP
app.config['MAIL_SERVER'] = 'smtp.qq.com'
app.config['MAIL_PORT'] = 465
app.config['MAIL_USE_SSL'] = True
app.config['MAIL_USE_TLS'] = False
app.config['MAIL_USERNAME'] = '1577418482@qq.com'
app.config['MAIL_PASSWORD'] = 'jtxqxfobatvlfgac'
app.config['MAIL_DEFAULT_SENDER'] = '1577418482@qq.com'
mail = Mail(app)

# 用于会话加密的密钥
app.secret_key = secrets.token_hex(16)

# 存储验证码的字典，格式为 {邮箱: (验证码, 过期时间)}
verification_codes = {}

# 全局变量存储拓扑结果
topology_results = {}

def get_db():
    """获取数据库连接"""
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
        db.row_factory = sqlite3.Row
    return db

def init_db():
    """初始化数据库（包含问答会话表和用户表）"""
    logger.info("开始初始化数据库...")
    with app.app_context():
        db_path = os.path.abspath(DATABASE)
        logger.info(f"数据库文件路径: {db_path}")
        
        # 确保数据库目录存在
        db_dir = os.path.dirname(db_path)
        if not os.path.exists(db_dir):
            os.makedirs(db_dir)
        
        # 如果数据库已存在，跳过初始化
        if os.path.exists(db_path):
            logger.info("数据库已存在，跳过初始化")
            return
        
        logger.info("数据库文件不存在，创建新数据库...")
        with closing(get_db()) as db:
            try:
                # 检查并添加缺少的列
                cursor = db.cursor()
                cursor.execute("PRAGMA table_info(topologies)")
                columns = [column[1] for column in cursor.fetchall()]
                
                if 'user_id' not in columns:
                    cursor.execute('ALTER TABLE topologies ADD COLUMN user_id TEXT DEFAULT "anonymous"')
                    logger.info("已添加 user_id 列到 topologies 表")
                
                schema = """
                CREATE TABLE IF NOT EXISTS topologies (
                    id TEXT PRIMARY KEY,
                    content TEXT NOT NULL,
                    max_nodes INTEGER DEFAULT 0,
                    created_at TEXT,
                    user_id TEXT DEFAULT 'anonymous'
                );
                
                CREATE TABLE IF NOT EXISTS nodes (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    label TEXT NOT NULL,
                    level INTEGER DEFAULT 0,
                    value REAL DEFAULT 1,
                    mastered INTEGER DEFAULT 0,
                    mastery_score REAL DEFAULT 0,
                    consecutive_correct INTEGER DEFAULT 0,
                    content_snippet TEXT,
                    PRIMARY KEY (topology_id, id),
                    FOREIGN KEY (topology_id) REFERENCES topologies (id)
                );
                
                CREATE TABLE IF NOT EXISTS edges (
                    topology_id TEXT,
                    from_node TEXT,
                    to_node TEXT,
                    label TEXT,
                    PRIMARY KEY (topology_id, from_node, to_node),
                    FOREIGN KEY (topology_id) REFERENCES topologies (id),
                    FOREIGN KEY (from_node) REFERENCES nodes (id),
                    FOREIGN KEY (to_node) REFERENCES nodes (id)
                );
                
                CREATE TABLE IF NOT EXISTS users (
                    id TEXT PRIMARY KEY,
                    username TEXT NOT NULL UNIQUE,
                    password TEXT NOT NULL,
                    email TEXT,
                    email_verified INTEGER DEFAULT 0,
                    created_at TEXT
                );
                
                -- 问答会话表
                CREATE TABLE IF NOT EXISTS quiz_sessions (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    node_id TEXT,
                    created_at TEXT,
                    consecutive_correct INTEGER DEFAULT 0,
                    mastered INTEGER DEFAULT 0,
                    FOREIGN KEY (topology_id) REFERENCES topologies (id),
                    FOREIGN KEY (node_id) REFERENCES nodes (id)
                );
                
                CREATE TABLE IF NOT EXISTS questions (
                    id TEXT PRIMARY KEY,
                    topology_id TEXT,
                    node_id TEXT,
                    question TEXT,
                    session_id TEXT,
                    created_at TEXT,
                    answered_at TEXT,
                    answer TEXT,
                    feedback TEXT,
                    correctness INTEGER DEFAULT 0,
                    FOREIGN KEY (topology_id) REFERENCES topologies (id),
                    FOREIGN KEY (node_id) REFERENCES nodes (id),
                    FOREIGN KEY (session_id) REFERENCES quiz_sessions (id)
                );
                
                CREATE TABLE IF NOT EXISTS password_resets (
                    id TEXT PRIMARY KEY,
                    user_id TEXT,
                    token TEXT,
                    created_at TEXT,
                    FOREIGN KEY (user_id) REFERENCES users (id)
                );
                """
                db.executescript(schema)
                db.commit()
                logger.info("数据库表创建成功")
            except Exception as e:
                logger.error(f"数据库初始化失败: {str(e)}", exc_info=True)
                raise

# 登录验证装饰器
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'username' not in session:
            if request.is_json:
                return jsonify({'status': 'error', 'message': '请先登录'}), 401
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# 验证登录函数
def verify_login(username, password):
    """验证登录"""
    with app.app_context():
        db = get_db()
        cursor = db.cursor()
        cursor.execute("SELECT password, email_verified, email FROM users WHERE username = ?", (username,))
        result = cursor.fetchone()
        
        if result and check_password_hash(result[0], password):
            # 检查邮箱是否已验证
            if result[1] == 0:  # email_verified = 0
                return False, "请先验证邮箱后再登录", result[2]  # 返回邮箱地址
            return True, "登录成功", None
    return False, "用户名或密码错误", None

# 注册新用户
def register_user(username, password, email=None):
    """注册新用户"""
    with app.app_context():
        db = get_db()
        cursor = db.cursor()
        
        # 检查用户名是否已存在
        cursor.execute("SELECT * FROM users WHERE username = ?", (username,))
        if cursor.fetchone():
            return False, "用户名已存在"
        
        # 检查邮箱是否已存在
        if email:
            cursor.execute("SELECT * FROM users WHERE email = ?", (email,))
            if cursor.fetchone():
                return False, "邮箱已存在"
        
        if not email:
            return False, "邮箱不能为空"
        
        try:
            hashed_password = generate_password_hash(password)
            cursor.execute("INSERT INTO users (username, password, email, email_verified) VALUES (?, ?, ?, ?)",
                         (username, hashed_password, email, 0))
            db.commit()
            
            # 自动发送验证邮件
            success, message = send_verification_email(email, username)
            if success:
                return True, "注册成功，验证邮件已发送至您的邮箱，请验证后登录"
            else:
                return True, f"注册成功，但发送验证邮件失败: {message}"
        except Exception as e:
            return False, f"注册失败: {str(e)}"

# 发送验证邮件
def send_verification_email(email, username):
    """发送验证邮件"""
    # 生成6位随机验证码
    code = ''.join(random.choices(string.digits, k=6))
    # 设置验证码有效期为10分钟
    expiration = datetime.now() + timedelta(minutes=10)
    verification_codes[email] = (code, expiration)
    
    # 发送验证码邮件
    try:
        msg = Message("邮箱验证码", recipients=[email])
        msg.body = f"亲爱的 {username}，\n\n您的邮箱验证码是：{code}，有效期10分钟。\n\n请在验证页面输入此验证码以完成邮箱验证。\n\n如果这不是您的操作，请忽略此邮件。"
        mail.send(msg)
        return True, "验证码已发送"
    except Exception as e:
        logger.error(f"发送验证邮件失败: {str(e)}", exc_info=True)
        return False, f"发送失败: {str(e)}"

def clean_json_string(s: str) -> str:
    """清洗模型输出，去除Markdown代码块标记"""
    s = re.sub(r"```(?:json)?", "", s)
    return s.strip()

# 增强JSON格式，处理各种复杂格式问题
def enhance_json_format(json_str: str) -> str:
    """增强JSON格式，处理各种复杂格式问题"""
    import json
    import re
    
    logger.info(f"原始JSON内容: {json_str[:200]}...")
    
    # 1. 移除可能的前导/尾随非JSON字符
    json_str = re.sub(r'^.*?(\[.*\]).*$', r'\1', json_str, flags=re.DOTALL)
    
    # 2. 处理行首缩进和换行
    json_str = re.sub(r'\n\s*', ' ', json_str)
    
    # 3. 处理可能的单引号问题
    json_str = re.sub(r"'", '"', json_str)
    
    # 4. 处理常见的格式错误（尝试修复）
    try:
        # 尝试直接解析
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        logger.error(f"初始解析失败: {e}")
        
        # 错误定位
        error_pos = e.pos
        logger.error(f"错误位置: {error_pos}, 附近内容: {json_str[error_pos-20:error_pos+20]}")
        
        # 5. 智能修复常见错误
        # 5.1 处理缺少逗号的情况
        if e.msg == "Expecting ',' delimiter":
            logger.info("尝试修复缺少逗号的问题...")
            json_list = list(json_str)
            
            # 在错误位置前查找可能缺少逗号的位置
            search_start = max(0, error_pos - 100)
            bracket_count = 0
            for i in range(error_pos-1, search_start, -1):
                if json_list[i] == ']':
                    bracket_count += 1
                elif json_list[i] == '[':
                    bracket_count -= 1
                
                # 找到合适的位置插入逗号
                if bracket_count == 0 and json_list[i] == ']':
                    json_list.insert(i+1, ',')
                    logger.info(f"在位置 {i+1} 插入逗号")
                    break
            
            # 尝试再次解析
            try:
                return json.loads(''.join(json_list))
            except json.JSONDecodeError as e2:
                logger.error(f"修复后仍失败: {e2}")
        
        # 5.2 处理未闭合的引号
        if e.msg.startswith('Expecting property name enclosed in double quotes'):
            logger.info("尝试修复未闭合的引号...")
            # 查找最近的未闭合引号并添加
            unclosed_quote_pos = json_str.rfind('"', 0, error_pos)
            if unclosed_quote_pos != -1:
                json_list = list(json_str)
                json_list.insert(error_pos, '"')
                logger.info(f"在位置 {error_pos} 添加引号")
                try:
                    return json.loads(''.join(json_list))
                except json.JSONDecodeError as e3:
                    logger.error(f"修复后仍失败: {e3}")
        
        # 6. 使用更宽松的解析器（作为最后的手段）
        try:
            import ast
            logger.info("尝试使用ast.literal_eval进行宽松解析...")
            parsed = ast.literal_eval(json_str)
            # 转换为标准JSON
            return json.loads(json.dumps(parsed))
        except Exception as e4:
            logger.error(f"宽松解析失败: {e4}")
            logger.error(f"无法修复的JSON内容: {json_str}")
            raise RuntimeError("无法解析API返回的JSON格式") from e4

def sanitize_text(text: str) -> str:
    """清理文本中的特殊字符，防止破坏JSON解析"""
    # 移除可能干扰JSON解析的字符
    text = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', text)  # 移除控制字符
    # 转义特殊字符
    text = text.replace('\\', '\\\\')
    text = text.replace('"', '\\"')
    return text

def extract_knowledge_from_text(text: str, max_nodes: int = 0, max_retries: int = MAX_RETRIES) -> list:
    """调用DeepSeek API提取适合树形结构的知识点层级关系"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 清理文本
    sanitized_text = sanitize_text(text)
    
    # 根据节点数量限制调整提示
    node_limit_prompt = ""
    if max_nodes > 0:
        node_limit_prompt = f"请确保最终提取的知识点数量不超过{max_nodes}个。"
    
    messages = [
        {"role": "system", "content": f"""你是一个知识图谱构建专家，能够从文本中提取知识点并构建树形结构。
请分析文本内容，识别出主要知识点及其层级关系（如父节点-子节点关系），
以JSON数组形式输出，每个元素格式为 [父知识点, 关系, 子知识点]。
关系应体现层级结构，如"包含"、"属于"、"是子类"等。确保输出格式正确，仅返回JSON数组。
{node_limit_prompt}"""},
        {"role": "user", "content": f"请从下面文本中提取知识点及其层级关系，输出JSON数组，每个元素格式为 [父知识点, 关系, 子知识点]：\n{sanitized_text}"}
    ]
    
    backoff = 2
    for attempt in range(1, max_retries + 1):
        try:
            logger.info(f"第{attempt}次尝试调用DeepSeek API...")
            
            # 使用OpenAI SDK调用API
            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=messages,
                max_tokens=1500
            )
            
            raw = response.choices[0].message.content
            cleaned = clean_json_string(raw or "")
            logger.info(f"API返回知识关系: {cleaned[:200]}...")
            
            # 保存原始响应用于调试
            # with open(f"api_response_{time.time()}.txt", "w", encoding="utf-8") as f:
            #     f.write(cleaned)
            
            # 增强JSON解析
            knowledge_edges = enhance_json_format(cleaned)
            
            # 验证输出格式
            if not isinstance(knowledge_edges, list):
                raise ValueError(f"API返回非数组格式: {type(knowledge_edges)}")
            for idx, item in enumerate(knowledge_edges):
                if not (isinstance(item, list) and len(item) == 3):
                    raise ValueError(f"API返回元素格式错误，应为三元组，位置 {idx}: {item}")
                    
            logger.info(f"成功解析知识关系，共{len(knowledge_edges)}条")
            return knowledge_edges
                
        except Exception as e:
            logger.error(f"API请求错误: {str(e)}", exc_info=True)
            if attempt < max_retries:
                logger.info(f"准备第{attempt+1}次重试...")
                time.sleep(backoff)
                backoff *= BACKOFF_FACTOR
            else:
                raise
    
    raise RuntimeError("多次重试后仍无法获取有效响应")

def parse_document(file_path):
    """解析文档内容，返回文本（新增PPT支持）"""
    file_ext = os.path.splitext(file_path)[1].lower()
    logger.info(f"开始解析文档: {file_path}, 类型: {file_ext}")
    
    try:
        if file_ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        elif file_ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    text += page_text
                    if page_num % 10 == 0:
                        logger.info(f"已解析PDF第 {page_num} 页")
            return text
        elif file_ext in ['.docx', '.doc']:
            doc = Document(file_path)
            full_text = []
            for para_num, para in enumerate(doc.paragraphs):
                full_text.append(para.text)
                if para_num % 50 == 0:
                    logger.info(f"已解析Word第 {para_num} 段落")
            return '\n'.join(full_text)
        elif file_ext == '.html':
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            soup = BeautifulSoup(html_content, 'lxml')
            text = soup.get_text()
            return ' '.join(text.split())
        elif file_ext in ['.ppt', '.pptx']:  # 新增PPT/PPTX支持
            text = ""
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if slide_num % 10 == 0:
                    logger.info(f"已解析PPT第 {slide_num} 页")
            return text
        else:
            logger.error(f"不支持的文件格式: {file_ext}")
            return None
    except Exception as e:
        logger.error(f"解析文档出错: {str(e)}", exc_info=True)
        return None

def extract_content_snippet(content: str, topic: str) -> str:
    """从原文中提取与主题相关的片段"""
    # 查找主题在内容中的位置
    index = content.lower().find(topic.lower())
    if index == -1:
        return ""

    # 提取上下文片段（主题前后各200个字符）
    start = max(0, index - 200)
    end = min(len(content), index + len(topic) + 200)
    snippet = content[start:end]
    
    # 确保片段包含主题
    if topic.lower() not in snippet.lower():
        return ""
    
    # 添加省略号表示截断
    if start > 0:
        snippet = "..." + snippet
    if end < len(content):
        snippet = snippet + "..."
    
    return snippet

def build_tree_structure(knowledge_edges, topology_id, content: str, max_nodes: int = 0, user_id=None):
    """构建树形知识图数据结构，保存原文片段并恢复掌握状态"""
    nodes = {}
    edges = []
    all_node_ids = set()
    
    # 收集所有节点
    for src, rel, tgt in knowledge_edges:
        all_node_ids.add(src)
        all_node_ids.add(tgt)
        
        # 确保节点存在
        if src not in nodes:
            # 提取原文片段
            snippet = extract_content_snippet(content, src)
            nodes[src] = {
                "id": src,
                "label": src,
                "title": src,  # 鼠标悬停时显示的完整标题
                "level": 0,    # 默认层级
                "value": 1,    # 节点大小基准
                "mastered": False,  # 知识点掌握状态
                "mastery_score": 0,  # 掌握分数
                "consecutive_correct": 0,  # 连续正确回答次数
                "content_snippet": snippet  # 保存原文片段
            }
        if tgt not in nodes:
            snippet = extract_content_snippet(content, tgt)
            nodes[tgt] = {
                "id": tgt,
                "label": tgt,
                "title": tgt,
                "level": 0,
                "value": 1,
                "mastered": False,
                "mastery_score": 0,
                "consecutive_correct": 0,
                "content_snippet": snippet
            }
        
        # 添加边
        edges.append({
            "from": src,
            "to": tgt,
            "label": rel,
            "title": rel,  # 鼠标悬停时显示的边关系
            "arrows": "to",
            "font": {
                "align": "middle"
            }
        })
    
    # 计算节点层级
    def calculate_level(node_id, current_level=0, visited=None):
        if visited is None:
            visited = set()
        if node_id in visited:
            return
        visited.add(node_id)
        if node_id in nodes:
            nodes[node_id]["level"] = max(nodes[node_id]["level"], current_level)
            # 递归设置子节点层级
            for edge in edges:
                if edge["from"] == node_id:
                    calculate_level(edge["to"], current_level + 1, visited)
    
    # 找到根节点（没有父节点的节点）
    root_candidates = all_node_ids.copy()
    for _, _, tgt in knowledge_edges:
        root_candidates.discard(tgt)
    
    root = next(iter(root_candidates)) if root_candidates else (list(nodes.keys())[0] if nodes else None)
    
    # 从根节点开始计算层级
    if root:
        calculate_level(root)
    
    # 计算节点重要性（连接数）
    for node_id in nodes:
        in_connections = sum(1 for edge in edges if edge["to"] == node_id)
        out_connections = sum(1 for edge in edges if edge["from"] == node_id)
        nodes[node_id]["value"] = max(1, in_connections + out_connections)
        
        # 从数据库获取并恢复节点的掌握状态
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT mastered, mastery_score, consecutive_correct FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node_status = cursor.fetchone()
            if node_status:
                nodes[node_id]["mastered"] = bool(node_status["mastered"])
                nodes[node_id]["mastery_score"] = node_status["mastery_score"]
                nodes[node_id]["consecutive_correct"] = node_status["consecutive_correct"]
    
    # 保存节点和边到数据库
    save_to_database(topology_id, list(nodes.values()), edges, content, max_nodes, user_id)
    
    # 转换为节点列表
    tree_nodes = list(nodes.values())
    
    return {
        "nodes": tree_nodes,
        "edges": edges,
        "root": root
    }

def save_to_database(topology_id, nodes, edges, content: str, max_nodes=0, user_id=None):
    """将知识图谱数据保存到数据库（保存原文内容和节点数量限制，关联用户）"""
    with app.app_context():
        db = get_db()
        cursor = db.cursor()
        
        # 获取当前用户ID，如果没有传入则使用默认值
        if user_id is None:
            user_id = 'anonymous'
        
        try:
            # 检查topologies表是否有user_id列，如果没有则添加
            cursor.execute("PRAGMA table_info(topologies)")
            columns = [column[1] for column in cursor.fetchall()]
            
            if 'user_id' not in columns:
                cursor.execute('ALTER TABLE topologies ADD COLUMN user_id TEXT DEFAULT "anonymous"')
                logger.info("已添加 user_id 列到 topologies 表")
                db.commit()
            
            # 保存拓扑图信息（包含原文内容、节点数量限制和用户ID）
            cursor.execute(
                "INSERT OR REPLACE INTO topologies (id, content, max_nodes, created_at, user_id) VALUES (?, ?, ?, ?, ?)",
                (topology_id, content, max_nodes, time.strftime('%Y-%m-%d %H:%M:%S'), user_id)
            )
            
            # 保存节点
            for node in nodes:
                cursor.execute(
                    """INSERT OR REPLACE INTO nodes 
                    (topology_id, id, label, level, value, mastered, mastery_score, consecutive_correct, content_snippet) 
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                    (topology_id, node["id"], node["label"], node["level"], node["value"], 
                     int(node["mastered"]), node["mastery_score"], node["consecutive_correct"], 
                     node.get("content_snippet", ""))
                )
            
            # 保存边
            for edge in edges:
                cursor.execute(
                    "INSERT OR REPLACE INTO edges (topology_id, from_node, to_node, label) VALUES (?, ?, ?, ?)",
                    (topology_id, edge["from"], edge["to"], edge["label"])
                )
            
            db.commit()
            logger.info(f"知识图谱 {topology_id} 保存成功，用户: {user_id}")
            
        except Exception as e:
            db.rollback()
            logger.error(f"保存知识图谱到数据库失败: {str(e)}")
            raise


def process_document(file_path, topology_id, max_nodes=0, user_id=None):
    """处理文档并生成树形知识图（支持节点数量限制）"""
    start_time = time.time()
    logger.info(f"开始处理文档: {file_path}, 拓扑ID: {topology_id}, 最大节点数: {max_nodes}")
    
    # 更新状态为处理中
    topology_results[topology_id] = {
        "status": "processing",
        "progress": 0,
        "message": "开始处理文档...",
        "created_at": time.strftime('%Y-%m-%d %H:%M:%S')
    }
    
    try:
        with app.app_context():
            update_progress(topology_id, 10, "解析文档内容...")
            text = parse_document(file_path)
            text = parse_document(file_path)

            if not text:
                topology_results[topology_id] = {
                    "status": "error",
                    "message": "无法解析文档内容"
                }
                logger.error(f"文档解析失败: {file_path}")
                return
            # ✅ 在这里添加缓存
            uploaded_documents[topology_id] = text
            
            # 新增：将全文内容存入数据库，便于后续问答检索
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "INSERT OR REPLACE INTO topologies (id, content, created_at) VALUES (?, ?, ?)",
                (topology_id, text, time.strftime('%Y-%m-%d %H:%M:%S'))
            )
            db.commit()
            
            update_progress(topology_id, 20, "准备提取知识层级...")
            text_length = len(text)
            if text_length < 100:
                topology_results[topology_id] = {
                    "status": "error",
                    "message": "文档内容过短，无法提取知识"
                }
                logger.warning(f"文档内容过短: {file_path}, 长度: {text_length}")
                return

            update_progress(topology_id, 60, "调用DeepSeek API提取知识层级...")
            knowledge_edges = extract_knowledge_from_text(text, max_nodes)
            logger.info(f"成功提取{len(knowledge_edges)}条知识层级关系")
            
            update_progress(topology_id, 80, "构建树形知识图并提取原文片段...")
            knowledge_graph = build_tree_structure(knowledge_edges, topology_id, text, max_nodes, user_id)
            
            processing_time = time.time() - start_time
            logger.info(f"树形知识图生成完成，耗时: {processing_time:.2f} 秒")
            
            update_progress(topology_id, 100, "处理完成")
            time.sleep(1)  # 确保前端有足够时间更新进度
            
            topology_results[topology_id] = {
                "status": "completed",
                "data": knowledge_graph,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph["nodes"]),
                "edge_count": len(knowledge_graph["edges"]),
                "processing_time": round(processing_time, 2),
                "text_length": text_length,
                "max_nodes": max_nodes  # 保存节点数量限制
            }
            
    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        logger.error(f"处理文档出错: {str(e)}", exc_info=True)
        topology_results[topology_id] = {
            "status": "error",
            "message": f"处理过程中出错: {str(e)}"
        }

def update_progress(topology_id, progress, message):
    """更新处理进度"""
    if topology_id in topology_results:
        topology_results[topology_id].update({
            'progress': progress,
            'message': message
        })
        logger.info(f"拓扑ID: {topology_id}, 进度: {progress}%, 消息: {message}")

@app.route('/api/generate', methods=['POST'])
@login_required
def generate_knowledge_graph():
    """处理用户点击生成按钮后的请求，支持节点数量控制"""
    if 'file' not in request.files:
        logger.error("文件上传错误: 没有文件")
        return jsonify({'status': 'error', 'message': '没有文件上传'}), 400
    
    file = request.files['file']
    if file.filename == '':
        logger.error("文件上传错误: 未选择文件")
        return jsonify({'status': 'error', 'message': '未选择文件'}), 400
    
    # 获取节点数量 - 确保从表单获取
    max_nodes = request.form.get('max_nodes', 0, type=int)
    
    # 检查文件大小
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    
    if file_size > 50 * 1024 * 1024:  # 50MB限制
        logger.error(f"文件上传错误: 文件过大 ({file_size/1024/1024:.2f} MB)")
        return jsonify({'status': 'error', 'message': '文件大小超过50MB限制'}), 400
    
    topology_id = str(uuid.uuid4())
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{topology_id}_{file.filename}")
    file.save(file_path)
    
    # 获取当前用户ID
    user_id = session.get('username')
    
    logger.info(f"文件上传成功: {file_path}, 大小: {file_size/1024/1024:.2f} MB, 最大节点数: {max_nodes}")
    
    # 启动处理线程，并在应用上下文中执行
    threading.Thread(
        target=lambda: with_app_context(process_document, file_path, topology_id, max_nodes, user_id)
    ).start()
    
    return jsonify({
        'status': 'success',
        'topology_id': topology_id,
        'message': '文档上传成功，正在生成知识图谱',
        'max_nodes': max_nodes  # 返回节点数量限制
    })

def with_app_context(func, *args, **kwargs):
    """在应用上下文中执行函数"""
    with app.app_context():
        func(*args, **kwargs)

@app.route('/api/topology/<topology_id>', methods=['GET'])
def get_topology(topology_id):
    if topology_id not in topology_results:
        # 尝试从数据库获取
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT id, content, max_nodes, created_at FROM topologies WHERE id = ?",
                (topology_id,)
            )
            topology = cursor.fetchone()
            
            if not topology:
                logger.error(f"获取拓扑图错误: ID不存在 ({topology_id})")
                return jsonify({'status': 'error', 'message': '拓扑图不存在'}), 404
            
            # 从数据库获取节点和边
            cursor.execute(
                "SELECT id, label, level, value, mastered, mastery_score, consecutive_correct, content_snippet FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            nodes = [dict(row) for row in cursor.fetchall()]
            
            cursor.execute(
                "SELECT from_node, to_node, label FROM edges WHERE topology_id = ?",
                (topology_id,)
            )
            edges = [dict(row) for row in cursor.fetchall()]
            
            knowledge_graph = {
                "nodes": nodes,
                "edges": edges,
                "root": next((node["id"] for node in nodes if node["level"] == 0), nodes[0]["id"] if nodes else None)
            }
            
            return jsonify({
                'status': 'success',
                'data': knowledge_graph,
                'created_at': topology["created_at"],
                'node_count': len(nodes),
                'edge_count': len(edges),
                'text_length': len(topology["content"]),
                'max_nodes': topology["max_nodes"]  # 返回节点数量限制
            })
    
    topology = topology_results[topology_id]
    
    if topology['status'] == 'processing':
        return jsonify({
            'status': 'processing',
            'progress': topology.get('progress', 0),
            'message': topology.get('message', '正在处理中'),
            'max_nodes': topology.get('max_nodes', 0)  # 返回节点数量限制
        })
    
    if topology['status'] == 'error':
        logger.error(f"获取拓扑图错误: {topology.get('message', '未知错误')}")
        return jsonify({
            'status': 'error',
            'message': topology.get('message', '生成知识图时出错')
        }), 500
    
    return jsonify({
        'status': 'success',
        'data': topology['data'],
        'created_at': topology['created_at'],
        'node_count': topology['node_count'],
        'edge_count': topology['edge_count'],
        'processing_time': topology['processing_time'],
        'text_length': topology.get('text_length', 0),
        'max_nodes': topology.get('max_nodes', 0)  # 返回节点数量限制
    })

@app.route('/api/topology/<topology_id>/set_max_nodes', methods=['POST'])
def set_topology_max_nodes(topology_id):
    """更新拓扑图的节点数量设置"""
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        max_nodes = data.get('max_nodes', 0)
        
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 更新拓扑图的节点数量设置
            cursor.execute(
                "UPDATE topologies SET max_nodes = ? WHERE id = ?",
                (max_nodes, topology_id)
            )
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'message': '节点数量设置已更新',
                'max_nodes': max_nodes
            })
            
    except Exception as e:
        logger.error(f"设置节点数量错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"设置节点数量时出错: {str(e)}"
        }), 500

@app.route('/api/topology/<topology_id>/regenerate', methods=['POST'])
def regenerate_topology(topology_id):
    """重新生成知识图谱，使用用户输入的新节点数量"""
    try:
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 从请求中获取新的节点数量
            data = request.get_json()
            if data is None:
                return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
            max_nodes = data.get('max_nodes', 0)  # 从请求中获取新的节点数量
            
            # 从数据库获取原文内容
            cursor.execute(
                "SELECT content FROM topologies WHERE id = ?",
                (topology_id,)
            )
            topology = cursor.fetchone()
            
            if not topology:
                return jsonify({
                    'status': 'error',
                    'message': '拓扑图不存在'
                }), 404
            
            content = topology["content"]
            
            update_progress(topology_id, 30, "重新提取知识层级...")
            knowledge_edges = extract_knowledge_from_text(content, max_nodes)  # 使用新的节点数量
            logger.info(f"重新生成成功提取{len(knowledge_edges)}条知识层级关系")
            
            # 保存当前节点的掌握状态
            cursor.execute(
                "SELECT id, mastered, mastery_score, consecutive_correct FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            mastery_states = {row["id"]: dict(row) for row in cursor.fetchall()}
            
            update_progress(topology_id, 70, "重新构建树形知识图...")
            knowledge_graph = build_tree_structure(knowledge_edges, topology_id, content, max_nodes)  # 使用新的节点数量
            
            # 恢复节点的掌握状态
            for node in knowledge_graph["nodes"]:
                node_id = node["id"]
                if node_id in mastery_states:
                    state = mastery_states[node_id]
                    node["mastered"] = bool(state["mastered"])
                    node["mastery_score"] = state["mastery_score"]
                    node["consecutive_correct"] = state["consecutive_correct"]
                    
                    # 更新数据库中的节点状态
                    cursor.execute(
                        """UPDATE nodes SET 
                        mastered = ?, mastery_score = ?, consecutive_correct = ?
                        WHERE topology_id = ? AND id = ?""",
                        (int(state["mastered"]), state["mastery_score"], 
                         state["consecutive_correct"], topology_id, node_id)
                    )
            
            # 更新拓扑图的节点数量设置到数据库
            cursor.execute(
                "UPDATE topologies SET max_nodes = ? WHERE id = ?",
                (max_nodes, topology_id)
            )
            
            # 更新处理结果
            topology_results[topology_id] = {
                "status": "completed",
                "data": knowledge_graph,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph["nodes"]),
                "edge_count": len(knowledge_graph["edges"]),
                "processing_time": 0,
                "text_length": len(content),
                "max_nodes": max_nodes  # 保存新的节点数量限制
            }
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'message': '知识图谱重新生成成功',
                'node_count': len(knowledge_graph["nodes"]),
                'edge_count': len(knowledge_graph["edges"]),
                'max_nodes': max_nodes  # 返回新的节点数量限制
            })
            
    except Exception as e:
        logger.error(f"重新生成知识图谱错误: {str(e)}", exc_info=True)
        # 新增：检查数据库中是否已生成新图谱
        try:
            with app.app_context():
                db = get_db()
                cursor = db.cursor()
                cursor.execute(
                    "SELECT id FROM topologies WHERE id = ?",
                    (topology_id,)
                )
                topology = cursor.fetchone()
                if topology:
                    # 图谱已存在，返回成功
                    return jsonify({
                        'status': 'success',
                        'message': '知识图谱已生成（部分异常）',
                        'topology_id': topology_id
                    }), 200
        except Exception as e2:
            logger.error(f"异常处理时再次出错: {str(e2)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"重新生成知识图谱时出错: {str(e)}"
        }), 200


@app.route('/api/topology/<topology_id>/node/<node_id>/question', methods=['GET'])
def get_question(topology_id, node_id):
    """获取关于指定节点的问题（基于原文内容，支持会话管理）"""
    try:
        with app.app_context():
            # 从数据库获取节点信息
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            
            if not node:
                return jsonify({
                    'status': 'error',
                    'message': '节点不存在'
                }), 404
            
            node_label = node["label"]
            content_snippet = node["content_snippet"]
            
            # 检查是否已有活跃会话
            session_id = request.args.get('session_id')
            if session_id:
                cursor.execute(
                    "SELECT mastered, consecutive_correct FROM quiz_sessions WHERE id = ?",
                    (session_id,)
                )
                session = cursor.fetchone()
                if session and session["mastered"]:
                    return jsonify({
                        'status': 'success',
                        'mastered': True,
                        'message': '该知识点已掌握'
                    })
            
            # 创建或获取会话
            if not session_id:
                session_id = str(uuid.uuid4())
                cursor.execute(
                    "INSERT INTO quiz_sessions (id, topology_id, node_id, created_at) VALUES (?, ?, ?, ?)",
                    (session_id, topology_id, node_id, time.strftime('%Y-%m-%d %H:%M:%S'))
                )
                db.commit()
            
            # 获取会话状态
            cursor.execute(
                "SELECT consecutive_correct FROM quiz_sessions WHERE id = ?",
                (session_id,)
            )
            session = cursor.fetchone()
            consecutive_correct = session["consecutive_correct"] if session else 0
            
            # 生成问题（基于会话状态）
            question = generate_question(node_label, content_snippet, consecutive_correct)
            
            # 保存问题到数据库
            question_id = str(uuid.uuid4())
            cursor.execute(
                "INSERT INTO questions (id, topology_id, node_id, question, session_id, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                (question_id, topology_id, node_id, question, session_id, time.strftime('%Y-%m-%d %H:%M:%S'))
            )
            db.commit()
            
            return jsonify({
                'status': 'success',
                'data': {
                    'question_id': question_id,
                    'question': question,
                    'node_id': node_id,
                    'session_id': session_id
                }
            })
    except Exception as e:
        logger.error(f"获取问题错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"生成问题时出错: {str(e)}"
        }), 500

def generate_question(topic, context, consecutive_correct=0):
    """根据连续正确次数生成不同难度的问题"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    # 根据掌握程度生成不同难度的问题
    difficulty_map = {
        0: "基础概念题，用简洁的语言解释",
        1: "理解应用题，结合实例说明",
        2: "综合分析题，比较相关概念"
    }
    difficulty = difficulty_map.get(consecutive_correct, "进阶思考题，拓展应用场景")
    
    # 构建提示词
    messages = [
        {"role": "system", "content": f"""你是一个教育专家，能够基于原文内容生成有针对性的问题。
请生成一个{difficulty}，测试用户对"{topic}"的理解。
问题应该清晰明确，基于提供的原文内容，难度适合当前掌握程度。
只返回问题文本，不需要其他内容。"""},
        {"role": "user", "content": f"原文片段: {context}\n问题:"}
    ]
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=150
        )
        
        question = (response.choices[0].message.content or "").strip()
        return question
    except Exception as e:
        logger.error(f"生成问题出错: {str(e)}", exc_info=True)
        return f"关于{topic}的问题（基于原文）"

@app.route('/api/topology/<topology_id>/question/<question_id>/answer', methods=['POST'])
def answer_question(topology_id, question_id):
    """处理用户对问题的回答（支持会话状态管理）并更新节点状态"""
    try:
        with app.app_context():
            data = request.get_json()
            if data is None:
                return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
            answer = data.get('answer', '')
            node_id = data.get('node_id', '')
            session_id = data.get('session_id', '')
            
            if not answer or not node_id or not session_id:
                return jsonify({
                    'status': 'error',
                    'message': '缺少必要的参数'
                }), 400
            
            # 从数据库获取问题
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT question, node_id, session_id FROM questions WHERE id = ? AND topology_id = ?",
                (question_id, topology_id)
            )
            question_data = cursor.fetchone()
            
            if not question_data:
                return jsonify({
                    'status': 'error',
                    'message': '问题不存在'
                }), 404
            
            question = question_data["question"]
            stored_node_id = question_data["node_id"]
            stored_session_id = question_data["session_id"]
            
            if stored_node_id != node_id or stored_session_id != session_id:
                return jsonify({
                    'status': 'error',
                    'message': '问题与会话不匹配'
                }), 400
            
            # 从数据库获取节点信息
            cursor.execute(
                "SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            if not node:
                return jsonify({
                    'status': 'error',
                    'message': '节点不存在'
                }), 400
            
            node_label = node["label"]
            content_snippet = node["content_snippet"]
            
            # 调用DeepSeek评估回答
            evaluation = evaluate_answer(question, answer, node_label, content_snippet)
            
            # 确定回答是否正确
            is_correct = evaluation["correct"] if "correct" in evaluation else False
            feedback_text = evaluation["feedback"] if "feedback" in evaluation else "无法评估回答"
            
            # 更新问题状态
            cursor.execute(
                "UPDATE questions SET answered_at = ?, answer = ?, feedback = ?, correctness = ? WHERE id = ?",
                (time.strftime('%Y-%m-%d %H:%M:%S'), answer, feedback_text, 1 if is_correct else 0, question_id)
            )
            
            # 更新会话状态
            cursor.execute(
                "SELECT consecutive_correct, mastered FROM quiz_sessions WHERE id = ?",
                (session_id,)
            )
            session = cursor.fetchone()
            if not session:
                return jsonify({
                    'status': 'error',
                    'message': '问答会话不存在'
                }), 404
            
            current_consecutive = session["consecutive_correct"]
            current_mastered = session["mastered"]
            
            # 更新连续正确计数
            new_consecutive = current_consecutive + 1 if is_correct else 0
            new_mastered = 1 if new_consecutive >= 1 else 0  # 只需答对1次即可掌握
            
            cursor.execute(
                """UPDATE quiz_sessions SET 
                questions_answered = questions_answered + 1,
                consecutive_correct = ?,
                mastered = ?
                WHERE id = ?""",
                (new_consecutive, new_mastered, session_id)
            )
            
            # 更新节点的掌握状态
            cursor.execute(
                "SELECT mastery_score, consecutive_correct FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node_status = cursor.fetchone()
            current_node_score = node_status["mastery_score"] if node_status else 0
            current_node_consecutive = node_status["consecutive_correct"] if node_status else 0
            
            node_new_score = min(10, current_node_score + (1 if is_correct else -0.5))
            node_new_consecutive = new_consecutive
            node_new_mastered = new_mastered
            
            cursor.execute(
                "UPDATE nodes SET mastery_score = ?, consecutive_correct = ?, mastered = ? WHERE topology_id = ? AND id = ?",
                (node_new_score, node_new_consecutive, int(node_new_mastered), topology_id, node_id)
            )
            
            # 确保掌握状态正确更新（显式处理）
            if node_new_mastered:
                cursor.execute(
                    "UPDATE nodes SET mastered = 1 WHERE topology_id = ? AND id = ?",
                    (topology_id, node_id)
                )
            
            # 如果未掌握，生成下一个问题
            next_question = None
            next_question_id = None
            if not new_mastered:
                # 生成下一个问题（基于更新后的状态）
                next_question = generate_question(node_label, content_snippet, new_consecutive)
                if next_question:
                    next_question_id = str(uuid.uuid4())
                    cursor.execute(
                        "INSERT INTO questions (id, topology_id, node_id, question, session_id, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                        (next_question_id, topology_id, node_id, next_question, session_id, time.strftime('%Y-%m-%d %H:%M:%S'))
                    )
            
            db.commit()
            
            return jsonify({
                'status': 'success',
                'data': {
                    'correct': is_correct,
                    'feedback': feedback_text,
                    'mastered': new_mastered,
                    'consecutive_correct': new_consecutive,
                    'session_id': session_id,
                    'next_question': {
                        'id': next_question_id,
                        'question': next_question
                    } if next_question else None
                }
            })
    except Exception as e:
        logger.error(f"处理回答错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"处理回答时出错: {str(e)}"
        }), 500

def evaluate_answer(question, answer, topic, context):
    """调用DeepSeek评估回答是否正确（包含上下文）"""
    from openai import OpenAI
    
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url="https://api.deepseek.com"
    )
    
    messages = [
        {"role": "system", "content": """你是一个知识评估专家，能够准确判断用户回答的正确性，并给出详细反馈。
请评估用户对问题"{question}"的回答"{answer}"是否正确，参考原文片段："{context}"。
你的评估应该包括：
1. 判断回答是否正确（是/否）
2. 提供具体的反馈，解释为什么正确或错误
3. 如果回答错误，提供正确的信息
4. 如果回答正确，考虑是否需要进一步提问以深化理解

输出格式为JSON，包含以下字段：
{
  "correct": true/false,
  "feedback": "具体的反馈信息",
  "next_question": "如果需要进一步提问，这里是下一个问题，否则为null"
}""".replace("{question}", question).replace("{answer}", answer).replace("{context}", context)},
        {"role": "user", "content": f"问题: {question}\n回答: {answer}\n原文片段: {context}\n请评估这个回答是否正确，并提供反馈。如果正确，考虑是否需要进一步提问。"}
    ]
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=300
        )
        
        response_text = (response.choices[0].message.content or "").strip()
        logger.info(f"评估回答响应: {response_text[:200]}...")
        
        # 解析响应为JSON
        try:
            feedback = json.loads(clean_json_string(response_text))
            # 确保包含所有必要字段
            if 'correct' not in feedback:
                feedback['correct'] = False
            if 'feedback' not in feedback:
                feedback['feedback'] = "无法解析评估结果，请重试"
            if 'next_question' not in feedback:
                feedback['next_question'] = None
        except json.JSONDecodeError:
            # 如果解析失败，创建默认反馈
            feedback = {
                "correct": False,
                "feedback": "无法解析评估结果，请重试",
                "next_question": None
            }
        
        return feedback
    except Exception as e:
        logger.error(f"评估回答出错: {str(e)}", exc_info=True)
        return {
            "correct": False,
            "feedback": f"评估回答时出错: {str(e)}",
            "next_question": None
        }

@app.route('/api/topology/<topology_id>/ignore', methods=['POST'])
def ignore_nodes(topology_id):
    """忽略用户选择的节点"""
    try:
        with app.app_context():
            data = request.get_json()
            if data is None:
                return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
            ignored_nodes = data.get('ignored_nodes', [])
            
            db = get_db()
            cursor = db.cursor()
            
            # 获取所有节点
            cursor.execute(
                "SELECT id, label, level, value, mastered, mastery_score, content_snippet FROM nodes WHERE topology_id = ?",
                (topology_id,)
            )
            all_nodes = [dict(row) for row in cursor.fetchall()]
            
            # 获取所有边
            cursor.execute(
                "SELECT from_node, to_node, label FROM edges WHERE topology_id = ?",
                (topology_id,)
            )
            all_edges = [dict(row) for row in cursor.fetchall()]
            
            # 筛选节点（排除被忽略的）
            filtered_nodes = [node for node in all_nodes if node["id"] not in ignored_nodes]
            
            # 筛选边（只保留未被忽略节点之间的边）
            filtered_edges = []
            node_ids = set([node["id"] for node in filtered_nodes])
            for edge in all_edges:
                if edge["from_node"] in node_ids and edge["to_node"] in node_ids:
                    filtered_edges.append(edge)
            
            return jsonify({
                'status': 'success',
                'data': {
                    'nodes': filtered_nodes,
                    'edges': filtered_edges,
                    'root': next((node["id"] for node in filtered_nodes if node["level"] == 0), filtered_nodes[0]["id"] if filtered_nodes else None)
                }
            })
    except Exception as e:
        logger.error(f"忽略节点错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"忽略节点时出错: {str(e)}"
        }), 500

@app.route('/api/topology/<topology_id>/node/<node_id>/master', methods=['POST'])
def master_node(topology_id, node_id):
    """设置节点的掌握状态"""
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        mastered = data.get('mastered', False)
        
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            # 更新节点掌握状态
            cursor.execute(
                "UPDATE nodes SET mastered = ? WHERE topology_id = ? AND id = ?",
                (1 if mastered else 0, topology_id, node_id)
            )
            db.commit()
            
            # 获取节点信息返回
            cursor.execute(
                "SELECT id, label FROM nodes WHERE topology_id = ? AND id = ?",
                (topology_id, node_id)
            )
            node = cursor.fetchone()
            
            if not node:
                return jsonify({'status': 'error', 'message': '节点不存在'}), 404
            
            return jsonify({
                'status': 'success',
                'message': '节点状态更新成功',
                'node': dict(node)
            })
            
    except Exception as e:
        logger.error(f"更新节点状态错误: {str(e)}", exc_info=True)
        return jsonify({
            'status': 'error',
            'message': f"更新节点状态时出错: {str(e)}"
        }), 500

###智能助手模块
from openai import OpenAI
client = OpenAI(api_key=OPENAI_API_KEY, base_url="https://api.deepseek.com")

# 全局变量缓存上传文档内容，方便检索
uploaded_documents = {}  # topology_id: 原文全文字符串

def recommend_resources_based_on_question(question):
    """
    使用 DeepSeek API 根据用户问题推荐相关学习资源，返回格式：
    [
      {"title": "资源标题", "url": "链接", "snippet": "相关内容摘要"},
      ...
    ]
    """
    messages = [
        {"role": "system", "content": "你是一个学习资源推荐专家，能够根据用户的问题推荐最相关的高质量中文学习资料。请根据用户的问题，推荐5个高质量的可访问的中文学习资源，每个资源包含title、url、snippet，要求以JSON数组格式输出。只返回JSON数组，不要有多余解释。"},
        {"role": "user", "content": f"问题：{question}\n请推荐5个相关学习资源。"}
    ]
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=800
        )
        raw = response.choices[0].message.content.strip()
        # 尝试解析JSON
        try:
            # 处理可能的markdown代码块
            if raw.startswith('```json'):
                raw = raw[7:]
            if raw.endswith('```'):
                raw = raw[:-3]
            resources = json.loads(raw)
            if isinstance(resources, list):
                return resources
        except Exception:
            pass
        # 解析失败返回空
        return []
    except Exception as e:
        logger.error(f"学习资源推荐API错误: {str(e)}", exc_info=True)
        return []

@app.route('/api/chat', methods=['POST'])
def chat_with_knowledge():
    """
    用户交互问答接口：
    优先基于上传文档内容回答，若文档中没有相关内容，则调用网络智能问答。
    同时进行相关学习资源推荐，返回相关链接和内容片段。
    """
    try:
        data = request.json
        topology_id = data.get('topology_id', '')
        user_question = data.get('question', '').strip()
        
        if not user_question:
            return jsonify({'status': 'error', 'message': '问题不能为空'}), 400
        
        # 获取上传文档全文内容
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute("SELECT content FROM topologies WHERE id = ?", (topology_id,))
            row = cursor.fetchone()
            document_text = row["content"] if row else ""
        
        # 直接用DeepSeek API在文档内容中查找相关内容
        doc_search_prompt = (
            "你是一个文档检索助手。请在下方给定的文档内容中查找与用户问题最相关的原文片段，"
            "并直接用文档原文文本回答。回答时用Markdown格式对原文文字进行重新排版，可以更改与文本意思无关的序数词和特殊符号，不要改变原文有效文字，"
            "**所有数学公式必须用LaTeX语法，并用$...$（行内）或$$...$$（块级）包裹，且不要用Markdown代码块（```）或中括号[]包裹公式**，不要改变原文有效文字。如果文档中没有相关内容，请只回复'未找到'。\n"
            "文档内容：" + document_text + "\n用户问题：" + user_question + "\n请用文档原文回答："
        )
        messages = [
            {"role": "system", "content": "你是一个文档检索助手。"},
            {"role": "user", "content": doc_search_prompt}
        ]
        try:
            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=messages,
                max_tokens=512
            )
            doc_answer = response.choices[0].message.content.strip()
        except Exception as e:
            logger.error(f"DeepSeek文档检索API错误: {str(e)}", exc_info=True)
            doc_answer = ""
        
        deny_phrases = ["未找到", "没有相关内容", "查无相关", "未检索到", "未能找到", "未能检索到", "没有找到"]
        deny_matched = [deny for deny in deny_phrases if deny in doc_answer] if doc_answer else []
        logger.info(f"[问答调试] doc_answer: {doc_answer}")
        logger.info(f"[问答调试] deny_matched: {deny_matched}")
        if doc_answer and not deny_matched:
            answer = doc_answer  # 直接返回原始AI回答，保留Markdown
            source = "document"
        else:
            # 文档中查不到，调用智能网络问答接口
            answer = generate_answer_from_web(user_question)  # 直接返回原始AI回答，保留Markdown
            source = "web"
        # 推荐学习资源
        resources = recommend_resources_based_on_question(user_question)
        return jsonify({
            'status': 'success',
            'answer': answer,
            'resources': resources,
            'source': source
        })
        
    except Exception as e:
        logger.error(f"交互问答错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': f"交互问答出错: {str(e)}"}), 500

def generate_answer_from_web(question):
    """
    调用网络智能问答接口（DeepSeek）
    """
    messages = [
        {"role": "system", "content": (
            "你是一个智能助理，能够基于互联网资源回答各种问题。"
            "回答时用Markdown格式排版，不要添加不合理的换行，去除空白段落，所有数学公式必须用LaTeX语法。"
            "**行内公式请用$...$包裹，且必须与文字同行；独占一行或需要居中显示的公式必须与文字同行用$$...$$包裹。**"
            "禁止用markdown代码块（```）、中括号[]或其他符号包裹公式。"
            "只返回直接的答案内容，不要多余解释。"
        )},
        {"role": "user", "content": question}
    ]
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=messages,
            max_tokens=1024  # 增大输出长度
        )
        answer = response.choices[0].message.content.strip()
        logger.info(f"[问答调试] 网络AI原始返回内容: {repr(answer)}")
        return answer
    except Exception as e:
        logger.error(f"生成网络回答错误: {str(e)}", exc_info=True)
        return "抱歉，网络问答服务不可用。"

@app.teardown_appcontext
def close_db(exception):
    """关闭数据库连接"""
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()

###首页登录模块
@app.route('/api/register', methods=['POST'])
def api_register():
    """用户注册接口"""
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        email = data.get('email', '').strip()
        
        if not username or not password or not email:
            return jsonify({'status': 'error', 'message': '用户名、密码和邮箱不能为空'}), 400
        
        # 检查用户名或邮箱是否已被注册
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT id FROM users WHERE username = ? OR email = ?",
                (username, email)
            )
            existing_user = cursor.fetchone()
            if existing_user:
                return jsonify({'status': 'error', 'message': '用户名或邮箱已被注册'}), 400
        
        # 生成用户ID
        user_id = str(uuid.uuid4())
        # 哈希密码
        hashed_password = generate_password_hash(password)
        created_at = time.strftime('%Y-%m-%d %H:%M:%S')
        
        # 插入新用户
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO users (id, username, password, email, created_at) VALUES (?, ?, ?, ?, ?)",
                (user_id, username, hashed_password, email, created_at)
            )
            db.commit()
        
        return jsonify({'status': 'success', 'message': '注册成功，请登录'}), 201
    except Exception as e:
        logger.error(f"注册错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '注册时出错'}), 500

@app.route('/api/login', methods=['POST'])
def api_login():
    """用户登录接口"""
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        
        if not username or not password:
            return jsonify({'status': 'error', 'message': '用户名和密码不能为空'}), 400
        # 查询用户
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT id, password FROM users WHERE username = ?",
                (username,)
            )
            user = cursor.fetchone()
            
            if not user:
                return jsonify({'status': 'error', 'message': '用户不存在'}), 404
            # 验证密码
            if not check_password_hash(user["password"], password):
                return jsonify({'status': 'error', 'message': '密码错误'}), 401
            # 登录成功，创建会话
            session['user_id'] = user["id"]
            session['username'] = username
            session.permanent = True  # 使会话持久化
            
            return jsonify({'status': 'success', 'message': '登录成功'}), 200
    except Exception as e:
        logger.error(f"登录错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '登录时出错'}), 500

@app.route('/api/logout', methods=['POST'])
def api_logout():
    """用户登出接口"""
    try:
        session.clear()  # 清除会话
        return jsonify({'status': 'success', 'message': '登出成功'}), 200
    except Exception as e:
        logger.error(f"登出错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '登出时出错'}), 500

@app.route('/api/request_password_reset', methods=['POST'])
def request_password_reset():
    """请求密码重置链接接口"""
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        email = data.get('email', '').strip()
        
        if not email:
            return jsonify({'status': 'error', 'message': '邮箱不能为空'}), 400
        # 检查邮箱是否存在
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT id, username FROM users WHERE email = ?",
                (email,)
            )
            user = cursor.fetchone()
            
            if not user:
                return jsonify({'status': 'error', 'message': '邮箱未注册'}), 404
            # 生成密码重置令牌
            token = secrets.token_urlsafe(16)
            created_at = time.strftime('%Y-%m-%d %H:%M:%S')
            # 存储令牌和用户关联
            cursor.execute(
                "INSERT INTO password_resets (id, user_id, token, created_at) VALUES (?, ?, ?, ?)",
                (str(uuid.uuid4()), user["id"], token, created_at)
            )
            db.commit()
            # 发送密码重置邮件
            reset_url = f"http://localhost:5000/reset_password?token={token}"
            with app.app_context():
                msg = Message("密码重置请求", recipients=[email])
                msg.body = f"点击以下链接重置密码：{reset_url}\n如果不是您本人操作，请忽略此邮件。"
                mail.send(msg)
            
            return jsonify({'status': 'success', 'message': '密码重置链接已发送到您的邮箱'}), 200
    except Exception as e:
        logger.error(f"请求密码重置错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '请求密码重置时出错'}), 500

@app.route('/api/reset_password', methods=['POST'])
def reset_password():
    """重置密码接口"""
    try:
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        token = data.get('token', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not token or not new_password:
            return jsonify({'status': 'error', 'message': '令牌和新密码不能为空'}), 400
        # 验证令牌并获取用户
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT user_id FROM password_resets WHERE token = ?",
                (token,)
            )
            reset_request = cursor.fetchone()
            
            if not reset_request:
                return jsonify({'status': 'error', 'message': '无效或过期的令牌'}), 400
            
            user_id = reset_request["user_id"]
            # 更新用户密码
            hashed_password = generate_password_hash(new_password)
            cursor.execute(
                "UPDATE users SET password = ? WHERE id = ?",
                (hashed_password, user_id)
            )
            # 删除密码重置请求
            cursor.execute(
                "DELETE FROM password_resets WHERE token = ?",
                (token,)
            )
            
            db.commit()
            
            return jsonify({'status': 'success', 'message': '密码已重置，请登录'}), 200
    except Exception as e:
        logger.error(f"重置密码错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '重置密码时出错'}), 500

@app.route('/api/user', methods=['GET'])
def get_user_info():
    """获取当前登录用户信息"""
    try:
        if 'user_id' not in session:
            return jsonify({'status': 'error', 'message': '未登录'}), 401
        
        user_id = session['user_id']
        
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT id, username, email, created_at FROM users WHERE id = ?",
                (user_id,)
            )
            user = cursor.fetchone()
            
            if not user:
                return jsonify({'status': 'error', 'message': '用户不存在'}), 404
            
            return jsonify({
                'status': 'success',
                'data': {
                    'id': user['id'],
                    'username': user['username'],
                    'email': user['email'],
                    'created_at': user['created_at']
                }
            }), 200
    except Exception as e:
        logger.error(f"获取用户信息错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '获取用户信息时出错'}), 500

@app.route('/api/user', methods=['PUT'])
def update_user_info():
    """更新用户信息"""
    try:
        if 'user_id' not in session:
            return jsonify({'status': 'error', 'message': '未登录'}), 401
        
        user_id = session['user_id']
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        username = data.get('username', '').strip()
        email = data.get('email', '').strip()
        
        if not username or not email:
            return jsonify({'status': 'error', 'message': '用户名和邮箱不能为空'}), 400
        # 更新用户信息
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "UPDATE users SET username = ?, email = ? WHERE id = ?",
                (username, email, user_id)
            )
            db.commit()
            # 更新会话信息
            session['username'] = username
            
            return jsonify({'status': 'success', 'message': '用户信息已更新'}), 200
    except Exception as e:
        logger.error(f"更新用户信息错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '更新用户信息时出错'}), 500

@app.route('/api/user/password', methods=['PUT'])
def api_change_password():
    """修改用户密码"""
    try:
        if 'user_id' not in session:
            return jsonify({'status': 'error', 'message': '未登录'}), 401
        
        user_id = session['user_id']
        data = request.get_json()
        if data is None:
            return jsonify({'status': 'error', 'message': 'Invalid JSON'}), 400
        old_password = data.get('old_password', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not old_password or not new_password:
            return jsonify({'status': 'error', 'message': '旧密码和新密码不能为空'}), 400
        # 验证旧密码
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute(
                "SELECT password FROM users WHERE id = ?",
                (user_id,)
            )
            user = cursor.fetchone()
            
            if not user or not check_password_hash(user["password"], old_password):
                return jsonify({'status': 'error', 'message': '旧密码错误'}), 401
            # 更新新密码
            hashed_password = generate_password_hash(new_password)
            cursor.execute(
                "UPDATE users SET password = ? WHERE id = ?",
                (hashed_password, user_id)
            )
            db.commit()
            
            return jsonify({'status': 'success', 'message': '密码已修改，请重新登录'}), 200
    except Exception as e:
        logger.error(f"修改密码错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': '修改密码时出错'}), 500

@app.route('/api/topology/status/<topology_id>', methods=['GET'])
def get_topology_status(topology_id):
    """获取拓扑图处理状态"""
    if topology_id not in topology_results:
        return jsonify({
            'status': 'error',
            'message': '拓扑图不存在'
        }), 404
    
    topology = topology_results[topology_id]
    
    return jsonify({
        'status': 'success',
        'data': {
            'topology_id': topology_id,
            'status': topology['status'],
            'progress': topology.get('progress', 0),
            'message': topology.get('message', ''),
            'created_at': topology.get('created_at', ''),
            'node_count': topology.get('node_count', 0),
            'edge_count': topology.get('edge_count', 0),
            'processing_time': topology.get('processing_time', 0),
            'text_length': topology.get('text_length', 0),
            'max_nodes': topology.get('max_nodes', 0)
        }
    }), 200

@app.route('/api/topologies', methods=['GET'])
def get_topologies():
    """获取当前用户所有拓扑图列表"""
    try:
        user_id = session.get('username', 'anonymous')
        
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            
            cursor.execute("""
                SELECT id, created_at, max_nodes, 
                       substr(content, 1, 100) as content_preview
                FROM topologies 
                WHERE user_id = ? 
                ORDER BY created_at DESC
            """, (user_id,))
            
            topologies = []
            for row in cursor.fetchall():
                topologies.append({
                    'id': row['id'],
                    'created_at': row['created_at'],
                    'max_nodes': row['max_nodes'],
                    'content_preview': row['content_preview'] + '...' if len(row['content_preview']) >= 100 else row['content_preview']
                })
            
            return jsonify({
                'status': 'success',
                'topologies': topologies
            })
            
    except Exception as e:
        logger.error(f"获取拓扑图列表出错: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': f"获取拓扑图列表失败: {str(e)}"}), 500

@app.route('/')
def index():
    """主页"""
    # 检查用户是否已登录，未登录则跳转到登录页面
    if 'username' not in session:
        return redirect(url_for('login'))
    return render_template('index.html', username=session.get('username'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    """用户登录"""
    if request.method == 'POST':
        data = request.get_json() if request.is_json else request.form
        username = data.get('username')
        password = data.get('password')
        
        success, message, email = verify_login(username, password)
        if success:
            session['username'] = username
            if request.is_json:
                return jsonify({'status': 'success', 'message': '登录成功'})
            else:
                return redirect(url_for('index'))
        else:
            if "请先验证邮箱" in message and email:
                # 重新发送验证邮件
                send_verification_email(email, username)
                if request.is_json:
                    return jsonify({
                        'status': 'error', 
                        'message': f'{message}，验证邮件已重新发送',
                        'email': email,
                        'require_verification': True
                    }), 403
                else:
                    return render_template('login.html', 
                                        error=f'{message}，验证邮件已重新发送至 {email}',
                                        show_verification_link=True,
                                        email=email)
            else:
                if request.is_json:
                    return jsonify({'status': 'error', 'message': message}), 401
                else:
                    return render_template('login.html', error=message)
    
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    """用户注册"""
    if request.method == 'POST':
        data = request.get_json() if request.is_json else request.form
        username = data.get('username')
        password = data.get('password')
        confirm_password = data.get('confirm_password')
        email = data.get('email')
        
        if not username or not password or not email:
            message = '用户名、密码和邮箱都不能为空'
        elif password != confirm_password:
            message = '两次输入的密码不一致'
        elif not re.match(r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$', email):
            message = '请输入有效的邮箱地址'
        else:
            success, message = register_user(username, password, email)
            if success:
                if request.is_json:
                    return jsonify({
                        'status': 'success', 
                        'message': message,
                        'email': email
                    })
                else:
                    flash(message, 'success')
                    return render_template('verify_email.html', email=email, username=username)
        
        if request.is_json:
            return jsonify({'status': 'error', 'message': message}), 400
        else:
            return render_template('register.html', error=message)
    
    return render_template('register.html')

@app.route('/logout')
def logout():
    """用户登出"""
    session.pop('username', None)
    return redirect(url_for('login'))

@app.route('/dashboard')
@login_required
def dashboard():
    """用户仪表板"""
    from datetime import datetime
    current_time = datetime.now()
    return render_template('dashboard.html', 
                         username=session.get('username'),
                         current_time=current_time)

@app.route('/forgot_password', methods=['GET', 'POST'])
def forgot_password():
    """忘记密码"""
    if request.method == 'POST':
        email = request.form.get('email')
        if email:
            # 这里可以添加发送重置密码邮件的逻辑
            flash('密码重置链接已发送到您的邮箱', 'success')
        else:
            flash('请输入邮箱地址', 'error')
    
    return render_template('forgot_password.html')

@app.route('/verify_email', methods=['GET', 'POST'])
def verify_email():
    """邮箱验证"""
    if request.method == 'POST':
        data = request.get_json() if request.is_json else request.form
        email = data.get('email')
        code = data.get('code')
        
        if email in verification_codes:
            stored_code, expiration = verification_codes[email]
            if datetime.now() > expiration:
                message = '验证码已过期'
                if request.is_json:
                    return jsonify({'status': 'error', 'message': message}), 400
                else:
                    flash(message, 'error')
            elif code == stored_code:
                # 更新用户邮箱验证状态
                with app.app_context():
                    db = get_db()
                    cursor = db.cursor()
                    cursor.execute("UPDATE users SET email_verified = 1 WHERE email = ?", (email,))
                    db.commit()
                
                del verification_codes[email]
                message = '邮箱验证成功，现在可以登录了'
                if request.is_json:
                    return jsonify({'status': 'success', 'message': message})
                else:
                    flash(message, 'success')
                    return redirect(url_for('login'))
            else:
                message = '验证码错误'
                if request.is_json:
                    return jsonify({'status': 'error', 'message': message}), 400
                else:
                    flash(message, 'error')
        else:
            message = '验证码不存在或已过期'
            if request.is_json:
                return jsonify({'status': 'error', 'message': message}), 400
            else:
                flash(message, 'error')
    
    # GET请求或验证失败时
    email = request.args.get('email', '')
    username = request.args.get('username', '')
    return render_template('verify_email.html', email=email, username=username)

@app.route('/send_verification', methods=['POST'])
def send_verification():
    """发送验证码"""
    data = request.get_json() if request.is_json else request.form
    email = data.get('email')
    username = data.get('username', 'User')
    
    if email:
        success, message = send_verification_email(email, username)
        if request.is_json:
            if success:
                return jsonify({'status': 'success', 'message': message})
            else:
                return jsonify({'status': 'error', 'message': message}), 500
        else:
            if success:
                flash(message, 'success')
            else:
                flash(message, 'error')
    else:
        message = '请输入邮箱地址'
        if request.is_json:
            return jsonify({'status': 'error', 'message': message}), 400
        else:
            flash(message, 'error')
    
    # 重定向回邮箱验证页面
    return redirect(url_for('verify_email', email=email, username=username))

@app.route('/change_password', methods=['GET', 'POST'])
@login_required
def change_password():
    """修改密码"""
    if request.method == 'POST':
        old_password = request.form.get('old_password')
        new_password = request.form.get('new_password')
        confirm_password = request.form.get('confirm_password')
        username = session.get('username')
        
        if not old_password or not new_password or not confirm_password:
            flash('所有字段都不能为空', 'error')
        elif new_password != confirm_password:
            flash('两次输入的新密码不一致', 'error')
        elif not verify_login(username, old_password):
            flash('原密码错误', 'error')
        else:
            # 更新密码
            with app.app_context():
                db = get_db()
                cursor = db.cursor()
                new_hashed = generate_password_hash(new_password)
                cursor.execute("UPDATE users SET password = ? WHERE username = ?", (new_hashed, username))
                db.commit()
            
            flash('密码修改成功，请重新登录', 'success')
            session.clear()
            return redirect(url_for('login'))
    
    return render_template('change_password.html')

@app.route('/test_email')
def test_email():
    """测试邮件功能"""
    try:
        success, message = send_verification_email("test@example.com", "TestUser")
        return f"邮件测试结果: {message}"
    except Exception as e:
        return f"邮件测试失败: {str(e)}"

@app.route('/test_verification')
def test_verification():
    """测试邮箱验证功能的页面"""
    return render_template('test_verification.html')


if __name__ == '__main__':
    # 创建必要的文件夹
    folders = ['uploads', 'static/css', 'static/js', 'templates']
    for folder in folders:
        if not os.path.exists(folder):
            os.makedirs(folder)
    # 初始化数据库
    try:
        init_db()
        logger.info("数据库初始化完成")
    except Exception as e:
        logger.error(f"数据库初始化异常: {str(e)}", exc_info=True)
        logger.info("尝试继续运行，但可能会出现数据库相关错误")
    # 拓扑图处理结果存储（使用全局变量）
    topology_results = {}
    uploaded_documents = {}
    
    logger.info("智能助教系统启动中...")
    app.run(debug=True, port=5000)